#!/usr/bin/env python3
"""Generate realistic, fictional company Office files for every C197 lab."""

from datetime import date, datetime, timedelta
import csv
import glob
import importlib
import json
import math
import os
import random
import re
import sys

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor as DocRGB
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PIn, Pt as PPt

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C

NAVY = "172B4D"; BLUE = "1F6FEB"; TEAL = "10B981"; VIOLET = "7C3AED"
AMBER = "F59E0B"; RED = "DC2626"; GREY = "5B6372"; LINE = "DDE5EF"
LIGHT = "F5F8FC"; ICE = "EAF2FF"; MINT = "EAF8F3"; WHITE = "FFFFFF"


def repo_root(start):
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return d
    raise RuntimeError("Course repository not found")


def load_activities():
    acts = []
    for path in sorted(glob.glob(os.path.join(HERE, "data_domain[0-9]*.py"))):
        name = os.path.basename(path)[:-3]
        if not re.fullmatch(r"data_domain\d+", name):
            continue
        n = re.search(r"\d+", name).group()
        acts.extend(getattr(importlib.import_module(name), f"DOMAIN{n}"))
    return sorted(acts, key=lambda a: a["num"])


import time as _time
_START = _time.time()
REPO = repo_root(HERE)
ACTS = load_activities()


def lab_folder(a):
    path = os.path.join(REPO, "labs", f"lab-{a['num']:02d}-{C.LAB_SLUGS[a['num']]}")
    os.makedirs(os.path.join(path, "templates"), exist_ok=True)
    return path


def stem(a):
    return f"Lumina-Living-Lab-{a['num']:02d}"


def shade(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_doc_defaults(doc):
    sec = doc.sections[0]
    sec.top_margin = Inches(0.62); sec.bottom_margin = Inches(0.62)
    sec.left_margin = Inches(0.72); sec.right_margin = Inches(0.72)
    normal = doc.styles["Normal"]
    normal.font.name = "Arial"; normal.font.size = Pt(10.5); normal.font.color.rgb = DocRGB.from_string(NAVY)
    for name, size, color in [("Title", 28, NAVY), ("Heading 1", 17, BLUE), ("Heading 2", 13, NAVY), ("Heading 3", 11, TEAL)]:
        style = doc.styles[name]
        style.font.name = "Arial"; style.font.size = Pt(size); style.font.bold = True; style.font.color.rgb = DocRGB.from_string(color)


def add_doc_brand(doc, label):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    r = p.add_run(label.upper())
    r.bold = True; r.font.size = Pt(8.5); r.font.color.rgb = DocRGB.from_string(TEAL)


def callout(doc, title, body, fill=ICE):
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = table.cell(0, 0); shade(cell, fill); cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    p = cell.paragraphs[0]
    r = p.add_run(title + "  "); r.bold = True; r.font.color.rgb = DocRGB.from_string(BLUE)
    p.add_run(body)


def add_table(doc, headers, rows, widths=None):
    table = doc.add_table(rows=1, cols=len(headers)); table.style = "Table Grid"; table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]; cell.text = header; shade(cell, BLUE)
        for run in cell.paragraphs[0].runs:
            run.font.color.rgb = DocRGB(255, 255, 255); run.bold = True; run.font.size = Pt(9)
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = str(value)
            for run in cells[i].paragraphs[0].runs: run.font.size = Pt(8.8)
    if widths:
        for row in table.rows:
            for i, width in enumerate(widths): row.cells[i].width = Inches(width)
    return table



# Lab 04 needs real role data in the brief, otherwise every field comes back
# "need to check" and the learner has nothing to work with.
LAB04_ROLES = [
    ("Warehouse", "Picker Packer", 4, "2,300 - 2,600", "Jan 2027", "Ops Manager",
     "Peak season backlog; 6 leavers in the last year"),
    ("Warehouse", "Shift Supervisor", 1, "3,700 - 4,200", "Feb 2027", "Ops Manager",
     "One shift currently runs without a supervisor"),
    ("Online", "Customer Service", 3, "2,800 - 3,100", "Jan 2027", "Online Lead",
     "Reply times doubled after the last two resignations"),
    ("Retail", "Store Assistant", 2, "2,400 - 2,700", "Mar 2027", "Retail Manager",
     "Two new store openings confirmed for Q1"),
    ("Office", "Data Analyst", 1, "4,600 - 5,200", "Apr 2027", "Head of HR",
     "Requested to build people reporting; no current gap in service"),
]
LAB04_BUDGET = "FY2027 hiring budget: $38,000 additional monthly salary cost. Headcount cap: 10 new roles."



# Labs 1 and 2 use a two-part staff handbook.  The content is deliberately
# incomplete so "need to check" is a real outcome, not a formality.
HANDBOOK_P1 = [
    ("1. Who this handbook applies to",
     "This handbook applies to all permanent staff at Lumina Living Pte Ltd across the retail, "
     "online and warehouse teams. Part-time staff are covered except where a section says otherwise. "
     "Agency and contract staff are covered by their own agreements."),
    ("2. Notice period — what we currently do",
     "Staff who have passed probation give one month's written notice. Staff still on probation give "
     "one week. Notice starts the day after the written notice is received. The company gives the same "
     "notice it asks for. Pay in lieu of notice has been agreed case by case in the past; the handbook "
     "does not currently set a rule for it."),
    ("3. Probation — what we currently do",
     "New staff serve three months' probation. The hiring manager holds a review at six weeks and again "
     "before the end. Probation has been extended by one month in some cases. The handbook does not say "
     "who approves an extension or how many times it may happen."),
    ("4. Flexible working — what staff have asked for",
     "Office and online teams have asked for two days a week from home. Warehouse and retail roles are "
     "on site by nature. Exit interviews in the last year repeatedly named inflexible hours as a reason "
     "for leaving. No flexible working policy exists yet."),
    ("5. What we must be careful about",
     "Anything about notice, pay in lieu, dismissal or statutory leave has legal consequences and must "
     "be checked by a qualified adviser before it is published to staff."),
    ("6. Who approves changes",
     "The Head of HR approves handbook wording. Anything with a legal consequence also needs sign-off "
     "from the company's legal adviser before release."),
]

HANDBOOK_P2 = [
    ("1. Who this handbook applies to",
     "This is part two of the Lumina Living staff handbook and applies to the same staff as part one."),
    ("2. Probation — what we currently do",
     "New staff serve three months' probation with a review at six weeks. Managers have extended "
     "probation informally. There is no written rule on how long an extension may run, who signs it "
     "off, or what happens to notice periods during an extension."),
    ("3. Flexible working — what staff have asked for",
     "Two days a week from home for office and online roles. Warehouse and retail roles are on site. "
     "The company has not decided whether this becomes a right, a manager's discretion, or a trial."),
    ("4. Leave — what we currently do",
     "Annual leave is 14 days rising to 18 after three years. Carry-over of up to five days has been "
     "allowed in practice but is not written down. Statutory leave entitlements are not restated here."),
    ("5. What we must be careful about",
     "Anything about leave entitlement, notice or dismissal must be checked by a qualified adviser "
     "before publication."),
    ("6. Who approves changes",
     "Head of HR, with legal sign-off for anything carrying a legal consequence."),
]


def save_handbook(a, folder, sections, suffix):
    doc = Document(); set_doc_defaults(doc)
    add_doc_brand(doc, f"{C.COURSE_CODE} · Lab {a['num']} · Company source")
    doc.add_paragraph(C.COMPANY, style="Title")
    p = doc.add_paragraph("Staff Handbook — working draft")
    p.runs[0].bold = True; p.runs[0].font.size = Pt(18)
    p.runs[0].font.color.rgb = DocRGB.from_string(TEAL)
    p = doc.add_paragraph("FICTIONAL COMPANY TRAINING MATERIAL")
    p.runs[0].bold = True; p.runs[0].font.color.rgb = DocRGB.from_string(GREY)
    callout(doc, "Before you draft",
            "This draft is deliberately incomplete. Where it does not state a rule, say so rather than "
            "inventing one.", MINT)
    for heading, body in sections:
        doc.add_heading(heading, level=1)
        doc.add_paragraph(body)
    path = os.path.join(folder, f"{stem(a)}-{suffix}.docx")
    doc.save(path); print("Saved", path)



# Labs 3 and 4 share the same staff data: Lab 3 analyses it in Excel and saves
# the method; Lab 4 has the skill read it from Word.
Q1_STAFF = [
    ("Retail", "Store Assistant", 18, 2450, 4),
    ("Retail", "Store Supervisor", 6, 3800, 1),
    ("Online", "Customer Service", 12, 2900, 3),
    ("Online", "Content Executive", 5, 3600, 1),
    ("Warehouse", "Picker Packer", 22, 2300, 6),
    ("Warehouse", "Forklift Operator", 7, 2900, 2),
    ("Warehouse", "Shift Supervisor", 5, 3900, 1),
    ("Office", "Finance Executive", 6, 4200, 0),
    ("Office", "HR Executive", 4, 3900, 1),
    ("Office", "IT Support", 3, 4400, 0),
]


def build_q1_staff(a):
    wb = Workbook(); add_readme_sheet(wb, a)
    ws = wb.create_sheet("Staff_List")
    style_sheet(ws, {"A": 8, "B": 16, "C": 24, "D": 12, "E": 20, "F": 20})
    ws.append(["Ref", "Team", "Role", "Headcount",
               "Average monthly salary", "Left in last 12 months"])
    style_header(ws)
    for i, (team, role, head, salary, leavers) in enumerate(Q1_STAFF, 1):
        ws.append([f"R{i:02d}", team, role, head, salary, leavers])
        ws[f"E{i+1}"].number_format = "#,##0"
    ws.freeze_panes = "A2"
    return wb


def save_q1_update(a, folder):
    doc = Document(); set_doc_defaults(doc)
    add_doc_brand(doc, f"{C.COURSE_CODE} · Lab {a['num']} · Company source")
    doc.add_paragraph(C.COMPANY, style="Title")
    p = doc.add_paragraph("Q1 People Update — for the leadership team")
    p.runs[0].bold = True; p.runs[0].font.size = Pt(18)
    p.runs[0].font.color.rgb = DocRGB.from_string(TEAL)
    p = doc.add_paragraph("FICTIONAL COMPANY TRAINING MATERIAL")
    p.runs[0].bold = True; p.runs[0].font.color.rgb = DocRGB.from_string(GREY)
    callout(doc, "Still to complete",
            "The staff figures below are missing. They are in the Q1 staff workbook in this folder.", MINT)
    for heading, body in [
        ("1. Purpose", "This update tells the leadership team where Lumina Living stands on staffing "
                       "at the end of Q1, and which team needs attention first."),
        ("2. Staff numbers by team", "[figures to be added from the Q1 staff workbook]"),
        ("3. Where the pressure is", "[the team with the highest leaver rate, and why]"),
        ("4. What we are asking for", "A decision on whether to prioritise replacement hiring in the "
                                      "team under most pressure."),
        ("5. Who approves this update", "Head of HR, before it goes to the leadership team."),
    ]:
        doc.add_heading(heading, level=1)
        doc.add_paragraph(body)
    path = os.path.join(folder, f"{stem(a)}-Q1-Update.docx")
    doc.save(path); print("Saved", path)



# Labs 11-13 (Skills) need their own paired data.
Q_TEAMS = {
    "Q1": [("Retail", 24, 2600, 5), ("Online", 17, 3100, 4),
           ("Warehouse", 34, 2500, 9), ("Office", 13, 4200, 1)],
    "Q2": [("Retail", 26, 2650, 3), ("Online", 19, 3150, 2),
           ("Warehouse", 31, 2700, 5), ("Office", 14, 4250, 2)],
}


def build_team_quarter(a, quarter):
    wb = Workbook(); add_readme_sheet(wb, a)
    ws = wb.create_sheet("Staff_List")
    style_sheet(ws, {"A": 8, "B": 16, "C": 12, "D": 20, "E": 20})
    ws.append(["Ref", "Team", "Headcount", "Average monthly salary",
               "Left in last 12 months"])
    style_header(ws)
    for i, (team, head, salary, leavers) in enumerate(Q_TEAMS[quarter], 1):
        ws.append([f"T{i:02d}", team, head, salary, leavers])
        ws[f"D{i+1}"].number_format = "#,##0"
    ws.freeze_panes = "A2"
    return wb


def save_q2_update(a, folder):
    doc = Document(); set_doc_defaults(doc)
    add_doc_brand(doc, f"{C.COURSE_CODE} · Lab {a['num']} · Company source")
    doc.add_paragraph(C.COMPANY, style="Title")
    p = doc.add_paragraph("Q2 People Update — for the leadership team")
    p.runs[0].bold = True; p.runs[0].font.size = Pt(18)
    p.runs[0].font.color.rgb = DocRGB.from_string(TEAL)
    callout(doc, "Still to complete",
            "The team figures below are missing. They are in the Q2 staff workbook in this folder.", MINT)
    for h, b in [("1. Purpose", "Where Lumina Living stands on staffing at the end of Q2."),
                 ("2. Staff numbers by team", "[to be added from the Q2 workbook]"),
                 ("3. Where the pressure is", "[the team with the highest leaver rate, and why]"),
                 ("4. Who approves this update", "Head of HR.")]:
        doc.add_heading(h, level=1); doc.add_paragraph(b)
    path = os.path.join(folder, f"{stem(a)}-Q2-Update.docx")
    doc.save(path); print("Saved", path)


DECK_STANDARD = """# Deck design standard

Apply this to any Lumina Living leadership deck.

1. Every slide title must state the conclusion, not name a topic.
   Write "Warehouse turnover is twice the company average", not "Turnover".
2. One message per slide. If a slide makes two points, split it.
3. Put a source note under any figure, naming the file and sheet it came from.
4. Flag any figure you cannot trace, rather than presenting it.
5. Keep the existing slide master, layouts, fonts and colours exactly as they are.
6. Speaker notes: three short lines saying what to say, not a script.
"""


def save_draft_deck(a, folder):
    prs = Presentation()
    titles = [("Turnover", "Turnover was reviewed this quarter across all teams."),
              ("Headcount", "Headcount numbers for the four teams are shown here."),
              ("Warehouse", "The warehouse team was looked at in detail."),
              ("Salary costs", "Salary costs by team over the last two quarters."),
              ("Hiring", "Hiring activity during the quarter."),
              ("Next steps", "Some possible next steps for consideration.")]
    for t, body in titles:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = t
        slide.placeholders[1].text = body
    path = os.path.join(folder, f"{stem(a)}-Draft-Deck.pptx")
    prs.save(path); print("Saved", path)
    md = os.path.join(folder, "deck-design-standard.md")
    with open(md, "w") as fh:
        fh.write(DECK_STANDARD)
    print("Saved", md)



# Lab 06 analyses people, not products: headcount, staff cost and leavers by
# team and month, with the assumptions kept on their own sheet.
PEOPLE_TEAMS = ["Retail", "Online", "Warehouse", "Office"]
PEOPLE_BASE = {"Retail": (24, 2600, 5), "Online": (17, 3100, 4),
               "Warehouse": (34, 2500, 9), "Office": (13, 4200, 1)}


def build_people_numbers(a):
    wb = Workbook(); add_readme_sheet(wb, a)
    ws = wb.create_sheet("Staff_List", 1)
    style_sheet(ws, {"A": 8, "B": 10, "C": 16, "D": 12, "E": 20, "F": 20})
    ws.append(["Ref", "Month", "Team", "Headcount", "Average monthly salary", "Leavers"])
    style_header(ws)
    i = 1
    for m_i, month in enumerate(MONTHS[:6]):
        for team in PEOPLE_TEAMS:
            head, salary, yr_leavers = PEOPLE_BASE[team]
            ws.append([f"S{i:03d}", month, team,
                       head + (m_i % 3) - 1, salary + m_i * 25,
                       max(0, round(yr_leavers / 6) + (1 if m_i == 2 and team == "Warehouse" else 0))])
            ws[f"E{i+1}"].number_format = "#,##0"
            i += 1
    ws.freeze_panes = "A2"

    plan = wb.create_sheet("Plan")
    style_sheet(plan, {"A": 10, "B": 16, "C": 16, "D": 20})
    plan.append(["Month", "Team", "Planned headcount", "Planned monthly cost"])
    style_header(plan)
    for m_i, month in enumerate(MONTHS[:6]):
        for team in PEOPLE_TEAMS:
            head, salary, _ = PEOPLE_BASE[team]
            plan.append([month, team, head, head * salary])
            plan.cell(row=plan.max_row, column=4).number_format = "#,##0"

    asm = wb.create_sheet("Assumptions")
    style_sheet(asm, {"A": 34, "B": 14, "C": 58})
    asm.append(["Assumption", "Value", "What it means"]); style_header(asm)
    for row in [("Employer cost on top of salary", 0.17, "CPF and benefits, as a share of salary"),
                ("Cost to replace one leaver", 3500, "Advertising, agency and manager time"),
                ("Target leaver rate", 0.10, "Company target, annual, as a share of headcount")]:
        asm.append(list(row))

    wb.create_sheet("Analysis")
    add_review_log(wb, "Checks")
    return wb



# Lab 05 ships a ready-made HR policy skill.  Learners upload it rather than
# writing it, which is how a real team shares one house standard.
HR_POLICY_SKILL = """# Lumina Living — HR policy wording standard

Use this whenever you draft or rewrite HR policy wording for staff.

## Structure every section in three clearly labelled parts

1. **What the policy says** — the rule itself, in one or two short sentences.
2. **How it works day to day** — what a manager or staff member actually does.
3. **Needs legal review** — anything that must be checked by a qualified
   adviser before the wording is published.

## Rules you must follow

- Use only what the open document actually says. Never invent a date, an
  amount, a notice period or an entitlement.
- After each fact, name the heading you took it from, in brackets.
- Where the document says nothing, write **need to check** instead of guessing.
- Never state a legal conclusion yourself. Flag it for review instead.
- Plain English. Short sentences. No jargon a new joiner would not understand.
- Keep the document's existing heading styles exactly as they are.

## Before you finish

End with a short list headed **Still to confirm**, naming every point marked
need to check and every point flagged for legal review.
"""


def save_policy_skill(a, folder):
    path = os.path.join(folder, "hr-policy-standard.md")
    with open(path, "w") as fh:
        fh.write(HR_POLICY_SKILL)
    print("Saved", path)



# Lab 05 ships a folder of real HR policy documents.  Learners hand the folder
# to Claude and ask it to build a skill from how these are actually written --
# which is how a company's house style really exists: in its documents.
POLICY_DOCS = [
    ("Annual-Leave-Policy", "Annual Leave Policy", [
        ("What the policy says",
         "Permanent staff receive 14 days of paid annual leave each year, rising to 18 days after "
         "three years of continuous service. Leave is granted in the calendar year in which it is earned."),
        ("How it works day to day",
         "Staff apply through their manager at least five working days in advance. Managers confirm "
         "cover before approving. Retail and warehouse teams may not take leave during the December "
         "peak without the department head's agreement (see 'Peak periods')."),
        ("Needs legal review",
         "Carry-over of unused leave into the following year has been permitted in practice but is "
         "not written down. Statutory minimum entitlements are not restated here and must be "
         "confirmed with a qualified adviser before publication."),
        ("Still to confirm",
         "Maximum carry-over days; whether carry-over needs approval; treatment of unused leave on "
         "resignation."),
    ]),
    ("Notice-Period-Policy", "Notice Period Policy", [
        ("What the policy says",
         "Staff who have completed probation give one month's written notice. Staff still on "
         "probation give one week. The company gives the same notice it asks for."),
        ("How it works day to day",
         "Notice starts the day after written notice is received by the line manager. The manager "
         "informs HR the same day so the leaver process can begin (see 'Leaver checklist')."),
        ("Needs legal review",
         "Payment in lieu of notice has been agreed case by case. No rule is written down and any "
         "wording on this must be checked by a qualified adviser."),
        ("Still to confirm",
         "Whether pay in lieu is at the company's discretion or the employee's; notice during an "
         "extended probation."),
    ]),
    ("Probation-Policy", "Probation Policy", [
        ("What the policy says",
         "New staff serve three months of probation. Employment may be confirmed, extended or ended "
         "at the end of that period."),
        ("How it works day to day",
         "The hiring manager holds a review at six weeks and again before the end of the third month. "
         "Outcomes are recorded on the probation review form and sent to HR."),
        ("Needs legal review",
         "Probation has been extended informally in some cases. The length of an extension, who "
         "approves it, and how many times it may happen are not written down. Any wording on "
         "ending employment during probation must be checked by a qualified adviser."),
        ("Still to confirm",
         "Maximum extension length; approver for an extension; notice period during an extension."),
    ]),
]


def save_policy_library(a, folder):
    """Write the existing HR policies as PDFs in their own folder."""
    lib = os.path.join(folder, "hr-policy-library")
    os.makedirs(lib, exist_ok=True)
    made = []
    for slug, title, sections in POLICY_DOCS:
        doc = Document(); set_doc_defaults(doc)
        _letterhead(doc, title)
        p = doc.add_paragraph("FICTIONAL COMPANY TRAINING MATERIAL · approved wording")
        p.runs[0].bold = True; p.runs[0].font.color.rgb = DocRGB.from_string(GREY)
        for heading, body in sections:
            doc.add_heading(heading, level=1)
            doc.add_paragraph(body)
        tmp = os.path.join(lib, slug + ".docx")
        doc.save(tmp)
        made.append(tmp)
    # Convert each to PDF, then drop the intermediate DOCX.
    import subprocess, glob as _g, shutil as _sh, tempfile as _tf
    with _tf.TemporaryDirectory() as td:
        subprocess.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", td] + made,
                       check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        for pdf in _g.glob(os.path.join(td, "*.pdf")):
            _sh.move(pdf, os.path.join(lib, os.path.basename(pdf)))
    for tmp in made:
        os.remove(tmp)
    print("Saved", lib, f"({len(POLICY_DOCS)} policy PDFs)")



HOW_WE_WRITE = """# How Lumina Living writes HR policy

Read this together with the approved policies in this folder. The PDFs show
the house style in practice; this file states the rules behind it.

## Every section has the same four parts, in this order

1. **What the policy says** — the rule itself, in one or two short sentences.
2. **How it works day to day** — what a manager or a staff member actually does.
3. **Needs legal review** — anything that must be checked by a qualified adviser
   before the wording is published to staff.
4. **Still to confirm** — a short list of every open point, so nothing quietly
   goes out unresolved.

## Rules

- Use only what the source document actually says. Never invent a date, an
  amount, a notice period or an entitlement.
- Name the heading each fact came from, in brackets, so it can be checked.
- Where the source is silent, write **need to check**. Do not fill the gap.
- Never state a legal conclusion. Flag it under 'Needs legal review' instead.
- Cross-reference other policies by name in single quotes, as the approved
  policies do — for example (see 'Leaver checklist').
- Plain English. Short sentences. No jargon a new joiner would not understand.
- Keep the existing heading styles of the document you are writing into.

## Tone

Neutral and factual. The reader is a member of staff who needs to know where
they stand, not a lawyer and not a manager.
"""


def save_policy_howto(a, folder):
    path = os.path.join(folder, "how-we-write-hr-policy.md")
    with open(path, "w") as fh:
        fh.write(HOW_WE_WRITE)
    print("Saved", path)



def _letterhead(doc, title):
    """Lumina Living letterhead: name bar, address line, then the title."""
    bar = doc.add_table(rows=1, cols=1)
    cell = bar.cell(0, 0); shade(cell, NAVY)
    p = cell.paragraphs[0]
    r = p.add_run("LUMINA LIVING PTE LTD")
    r.bold = True; r.font.size = Pt(16); r.font.color.rgb = DocRGB(255, 255, 255)
    sub = doc.add_paragraph()
    sr = sub.add_run("18 Kallang Avenue, #06-12, Singapore 339410   ·   "
                     "people@luminaliving.example   ·   Co. Reg. 2019xxxxxK")
    sr.font.size = Pt(8); sr.font.color.rgb = DocRGB.from_string(GREY)
    doc.add_paragraph(title, style="Title")


def save_policy_template(a, folder):
    """A blank, letterheaded policy template with the four house sections."""
    os.makedirs(os.path.join(folder, "templates"), exist_ok=True)
    doc = Document(); set_doc_defaults(doc)
    _letterhead(doc, "[Policy name]")
    p = doc.add_paragraph()
    pr = p.add_run("Human Resources   ·   Version [x.x]   ·   Approved by [name], Head of HR   "
                   "·   Effective [date]")
    pr.font.size = Pt(9); pr.font.color.rgb = DocRGB.from_string(GREY)
    callout(doc, "Before you publish",
            "Every section below must be completed. Leave 'need to check' in place wherever the "
            "source material is silent — do not delete it to make the document look finished.", MINT)
    for heading, hint in [
        ("1. What the policy says",
         "[The rule itself, in one or two short sentences.]"),
        ("2. How it works day to day",
         "[What a manager or staff member actually does. Cross-reference other policies by name "
         "in single quotes.]"),
        ("3. Needs legal review",
         "[Anything that must be checked by a qualified adviser before this is published.]"),
        ("4. Still to confirm",
         "[Every open point, listed. Nothing goes to staff while this list has entries.]"),
    ]:
        doc.add_heading(heading, level=1)
        ph = doc.add_paragraph(hint)
        ph.runs[0].italic = True
        ph.runs[0].font.color.rgb = DocRGB.from_string(GREY)
    doc.add_heading("Approval", level=1)
    add_table(doc, ["Role", "Name", "Date"],
              [("Drafted by", "", ""), ("Reviewed by (HR)", "", ""),
               ("Legal review", "", ""), ("Approved by (Head of HR)", "", "")],
              [2.2, 3.0, 1.6])
    path = os.path.join(folder, "templates", "HR-Policy-Template.docx")
    doc.save(path); print("Saved", path)



# Lab 05 needs a source with REAL facts, otherwise the skill correctly refuses
# to draft and the lab dead-ends.  This is what an HR team would actually have:
# a consultation note with decisions taken, gaps still open, and staff feedback.
LAB05_SOURCE = [
    ("1. Why we are writing these policies",
     "Exit interviews over the last 12 months named inflexible hours as a reason for leaving in "
     "9 of 21 cases, concentrated in the online and office teams. The Head of HR asked for two "
     "policies to be drafted for the January staff handbook update: flexible working, and "
     "carry-over of unused annual leave."),
    ("2. Flexible working — what the management team has agreed",
     "Office and online staff who have completed probation may work from home up to two days a "
     "week. Retail and warehouse roles are on site by nature and are out of scope. Days must be "
     "agreed with the line manager in advance and recorded in the team calendar. Core hours when "
     "everyone must be contactable are 10am to 4pm. Managers may require attendance on a "
     "home-working day where there is a business need, giving at least 24 hours' notice."),
    ("3. Flexible working — what is still open",
     "Whether the two days become a contractual right or remain at the manager's discretion has "
     "not been decided. There is no agreed process for a staff member to appeal a refused request. "
     "The treatment of staff who move from an on-site role to an office role has not been discussed."),
    ("4. Leave carry-over — what the management team has agreed",
     "Staff may carry a maximum of five unused annual leave days into the following year. Carried "
     "days must be used by 31 March or they lapse. Carry-over needs the line manager's approval, "
     "recorded on the leave system before 31 December."),
    ("5. Leave carry-over — what is still open",
     "Whether carried days are paid out on resignation has not been decided and legal has not been "
     "consulted. Whether long-term sick leave changes the carry-over limit is unresolved. The "
     "interaction with statutory minimum entitlement has not been checked."),
    ("6. What staff said",
     "'I would stay if I could work from home two days a week.' — online team, exit interview. "
     "'Nobody can tell me whether my four unused days roll over.' — office team, engagement survey. "
     "'The rule seems to change depending on who you ask.' — warehouse team, engagement survey."),
    ("7. Controls",
     "Both policies must go to the company's legal adviser before publication. Nothing may be "
     "stated as a legal entitlement in the draft. Every open point above must be visible in the "
     "draft rather than resolved by whoever writes it. The Head of HR approves the final wording."),
]


def save_lab05_source(a, folder):
    doc = Document(); set_doc_defaults(doc)
    _letterhead(doc, "HR Policy Consultation Note")
    p = doc.add_paragraph("Human Resources   ·   FICTIONAL COMPANY TRAINING MATERIAL   ·   "
                          "for the January staff handbook update")
    p.runs[0].font.size = Pt(9); p.runs[0].font.color.rgb = DocRGB.from_string(GREY)
    callout(doc, "How to use this note",
            "This is the only source for the two policies you are about to draft. Where it says a "
            "point is still open, that point must stay open in your draft.", MINT)
    for heading, body in LAB05_SOURCE:
        doc.add_heading(heading, level=1)
        doc.add_paragraph(body)
    path = os.path.join(folder, f"{stem(a)}-HR-Consultation-Note.docx")
    doc.save(path); print("Saved", path)



def save_blank_deck(a, folder):
    """A genuinely blank deck: title slide only, company colours on the master."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Lumina Living"
    slide.placeholders[1].text = "Q1 People Update  ·  for the leadership team"
    path = os.path.join(folder, f"{stem(a)}-Blank-Deck.pptx")
    prs.save(path); print("Saved", path)


def save_lab07_numbers(a, folder):
    """Headcount workbook with two ready-made charts, so the deck has something
    real to import even before the learner builds their own."""
    wb = build_people_numbers(a)
    src = wb["Staff_List"]
    n = src.max_row

    # A summary block the charts can point at: one row per team.
    ana = wb["Analysis"]
    ana["A1"] = "Team"; ana["B1"] = "Headcount"; ana["C1"] = "Leavers"; ana["D1"] = "Leaver rate"
    for c in ("A1", "B1", "C1", "D1"):
        ana[c].font = Font(name="Arial", bold=True, color=WHITE)
        ana[c].fill = PatternFill("solid", fgColor=BLUE)
    for i, team in enumerate(PEOPLE_TEAMS, start=2):
        ana[f"A{i}"] = team
        ana[f"B{i}"] = f'=SUMIF(Staff_List!C2:C{n},A{i},Staff_List!D2:D{n})'
        ana[f"C{i}"] = f'=SUMIF(Staff_List!C2:C{n},A{i},Staff_List!F2:F{n})'
        ana[f"D{i}"] = f"=IF(B{i}=0,0,C{i}/B{i})"
        ana[f"D{i}"].number_format = "0.0%"
    ana.column_dimensions["A"].width = 16
    for col in ("B", "C", "D"):
        ana.column_dimensions[col].width = 14

    rows = len(PEOPLE_TEAMS) + 1
    bar = BarChart()
    bar.title = "Warehouse loses people fastest"
    bar.y_axis.title = "Leaver rate"
    bar.add_data(Reference(ana, min_col=4, min_row=1, max_row=rows), titles_from_data=True)
    bar.set_categories(Reference(ana, min_col=1, min_row=2, max_row=rows))
    bar.height, bar.width = 8, 16
    ana.add_chart(bar, "F2")

    head = BarChart()
    head.title = "Headcount by team"
    head.add_data(Reference(ana, min_col=2, min_row=1, max_row=rows), titles_from_data=True)
    head.set_categories(Reference(ana, min_col=1, min_row=2, max_row=rows))
    head.height, head.width = 8, 16
    ana.add_chart(head, "F20")

    path = os.path.join(folder, f"{stem(a)}-People-Numbers.xlsx")
    wb.save(path); print("Saved", path, "(2 charts on Analysis)")



SETUP_ROWS = [
    ("Office add-in", "Claude panel opens in Word", "Home > Add-ins > Claude (Windows) or Tools > Add-ins > Claude (Mac)"),
    ("Office add-in", "Claude panel opens in Excel", "Same add-in; check it appears here too"),
    ("Office add-in", "Claude panel opens in PowerPoint", "Same add-in; check it appears here too"),
    ("Office add-in", "Claude panel opens in Outlook", "Optional; used only for the Outlook step in Lab 8"),
    ("Claude Desktop", "Signed in to Claude Desktop", "Same account as the Office add-in"),
    ("Claude Desktop", "Desktop can read a folder you choose", "Plus button > Add files or photos > choose this lab folder"),
    ("Microsoft 365 connector", "Microsoft 365 appears in Connectors", "Claude Desktop > Customize > Connectors"),
    ("Microsoft 365 connector", "Connected with a work account", "Read the permission screen before accepting"),
    ("Microsoft 365 connector", "A test search returns something", "Ask Claude to find any file; an empty result is still a result"),
    ("Claude for Chrome", "Extension installed and pinned", "claude.com/claude-for-chrome; Chrome only"),
    ("Claude for Chrome", "Permission mode set to Manually approve", "In the side panel; never Skip all approvals"),
]


def save_setup_checklist(a, folder):
    wb = Workbook(); add_readme_sheet(wb, a)
    ws = wb.create_sheet("Setup", 1)
    style_sheet(ws, {"A": 24, "B": 40, "C": 58, "D": 20, "E": 34})
    ws.append(["Route", "What to check", "Where to find it", "Result", "Note"])
    style_header(ws)
    for route, check, where in SETUP_ROWS:
        ws.append([route, check, where, "", ""])
    dv = DataValidation(type="list",
                        formula1='"Ready,Not available,Needs IT approval,Not checked"',
                        allow_blank=True)
    ws.add_data_validation(dv)
    dv.add(f"D2:D{len(SETUP_ROWS) + 1}")
    for r in range(2, len(SETUP_ROWS) + 2):
        ws[f"C{r}"].alignment = Alignment(wrap_text=True, vertical="top")
        ws[f"E{r}"].alignment = Alignment(wrap_text=True, vertical="top")
        ws.row_dimensions[r].height = 28
    ws.freeze_panes = "A2"
    path = os.path.join(folder, f"{stem(a)}-Setup-Checklist.xlsx")
    wb.save(path); print("Saved", path)



# Lab 09 reads across a folder, so it needs a folder worth reading: several
# files, in different formats, that do not entirely agree with each other.
LAB09_FILES = [
    ("Q1-Headcount-Report", "Q1 Headcount Report", [
        ("Summary",
         "Total headcount at 31 March was 88, against a plan of 84. Retail is 2 over plan following "
         "the two new store openings. Warehouse is 4 over plan because replacement hiring for "
         "leavers overlapped with the peak-season temporary contracts."),
        ("Leavers",
         "21 people left in the 12 months to March. Warehouse accounted for 9, online for 4, retail "
         "for 6 and office for 2. The warehouse leaver rate of 26.5% is the highest in the company "
         "and more than double the office rate."),
        ("Cost",
         "Monthly staff cost at 31 March was $253,400 against a plan of $241,000. The variance is "
         "explained by the overlap described above and is expected to unwind by June."),
    ]),
    ("Exit-Interview-Themes", "Exit Interview Themes — 12 months to March", [
        ("Method",
         "21 exit interviews were held. Themes were counted where a reason was raised without "
         "prompting. Percentages are of the 21 interviews, not of all staff."),
        ("What people said",
         "Inflexible hours were raised in 9 interviews, concentrated in online and office roles. "
         "Pay was raised in 7, almost all warehouse. Lack of progression was raised in 5. Relations "
         "with a direct manager were raised in 3."),
        ("Note on warehouse",
         "Warehouse leavers cited pay far more often than hours. A flexible-working policy is "
         "unlikely to change the warehouse figure on its own."),
    ]),
    ("Hiring-Pipeline-Note", "Hiring Pipeline Note — April", [
        ("Open roles",
         "Ten roles were approved for FY2027: four warehouse pickers, one warehouse shift "
         "supervisor, three online customer service, and two retail store assistants. The office "
         "data analyst role was held back."),
        ("Progress",
         "Six offers have been accepted. The two retail roles have not been advertised yet because "
         "the store opening dates moved to May."),
        ("Risk",
         "Warehouse pickers are taking an average of 21 days to fill, against a company average of "
         "nine. If this continues the peak-season cover will be short."),
    ]),
]

LAB09_CSV = [
    ("Team", "Headcount", "Plan", "Leavers_12m", "Monthly_cost"),
    ("Retail", 26, 24, 6, 66300),
    ("Online", 19, 18, 4, 58900),
    ("Warehouse", 30, 26, 9, 79200),
    ("Office", 13, 16, 2, 49000),
]


def save_lab09_folder(a, folder):
    """A small HR folder as it really arrives: reports as PDFs, plus a CSV."""
    lib = os.path.join(folder, "hr-quarter-files")
    os.makedirs(lib, exist_ok=True)
    made = []
    for slug, title, sections in LAB09_FILES:
        doc = Document(); set_doc_defaults(doc)
        _letterhead(doc, title)
        p = doc.add_paragraph("Human Resources   ·   FICTIONAL COMPANY TRAINING MATERIAL")
        p.runs[0].font.size = Pt(9); p.runs[0].font.color.rgb = DocRGB.from_string(GREY)
        for heading, body in sections:
            doc.add_heading(heading, level=1)
            doc.add_paragraph(body)
        tmp = os.path.join(lib, slug + ".docx")
        doc.save(tmp)
        made.append(tmp)
    import subprocess, glob as _g, shutil as _sh, tempfile as _tf
    with _tf.TemporaryDirectory() as td:
        subprocess.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", td] + made,
                       check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        for pdf in _g.glob(os.path.join(td, "*.pdf")):
            _sh.move(pdf, os.path.join(lib, os.path.basename(pdf)))
    for tmp in made:
        os.remove(tmp)
    with open(os.path.join(lib, "team-numbers-march.csv"), "w", newline="") as fh:
        csv.writer(fh).writerows(LAB09_CSV)
    print("Saved", lib, f"({len(LAB09_FILES)} PDFs + 1 CSV)")



# Lab 10 (Cowork) works on a week of raw HR inputs: the things that land on an
# HR desk between Monday and Friday and have to become one routine.
LAB10_INBOX = [
    ("Mon", "New starter", "Rachel Sim starts in Online on 3 March. Laptop and access not yet requested."),
    ("Mon", "Leaver", "Marcus Tan, Warehouse, resigned. Last day 28 March. Exit interview not booked."),
    ("Tue", "Leave request", "Priya Raman requests 5 days in April. Team already has two people away."),
    ("Tue", "Probation", "Hafiz Osman's probation review is due 12 March. No review form submitted."),
    ("Wed", "New starter", "Terrence Wong starts in Office on 10 March. Induction not scheduled."),
    ("Wed", "Policy question", "Three staff asked whether unused leave carries over. No written answer yet."),
    ("Thu", "Leaver", "Nur Syafiqah, Retail, resigned. Last day 21 March. Replacement not advertised."),
    ("Thu", "Probation", "Cheryl Ang passed probation on 5 March. Confirmation letter not sent."),
    ("Fri", "Headcount", "Warehouse is 4 over plan. Finance asked why by Monday."),
    ("Fri", "Hiring", "Two retail roles still unadvertised; store opening moved to May."),
]


def save_lab10_week(a, folder):
    """A week of raw HR inputs plus an empty routine tracker."""
    wb = Workbook(); add_readme_sheet(wb, a)
    ws = wb.create_sheet("This_Week", 1)
    style_sheet(ws, {"A": 8, "B": 18, "C": 70, "D": 20, "E": 18, "F": 26})
    ws.append(["Day", "Type", "What came in", "What must happen", "Who owns it", "By when"])
    style_header(ws)
    for day, kind, detail in LAB10_INBOX:
        ws.append([day, kind, detail, "", "", ""])
    for r in range(2, len(LAB10_INBOX) + 2):
        for col in ("C", "D", "F"):
            ws[f"{col}{r}"].alignment = Alignment(wrap_text=True, vertical="top")
        ws.row_dimensions[r].height = 30
    ws.freeze_panes = "A2"

    routine = wb.create_sheet("Daily_Routine")
    style_sheet(routine, {"A": 12, "B": 46, "C": 30, "D": 22})
    routine.append(["When", "What to do", "Where the information is", "Who checks it"])
    style_header(routine)
    routine["A3"] = "Fill this in with Cowork. One row per step of the daily routine."
    path = os.path.join(folder, f"{stem(a)}-This-Week.xlsx")
    wb.save(path); print("Saved", path)


def save_company_brief(a, folder):
    doc = Document(); set_doc_defaults(doc); add_doc_brand(doc, f"{C.COURSE_CODE} · Lab {a['num']} · Company source")
    doc.add_paragraph(C.COMPANY, style="Title")
    p = doc.add_paragraph(a["title"]); p.runs[0].bold = True; p.runs[0].font.size = Pt(18); p.runs[0].font.color.rgb = DocRGB.from_string(TEAL)
    p = doc.add_paragraph("FICTIONAL COMPANY TRAINING MATERIAL · FY2027 PLANNING CYCLE")
    p.runs[0].bold = True; p.runs[0].font.color.rgb = DocRGB.from_string(GREY)
    callout(doc, "Decision required", a["case"]["decision"], MINT)
    doc.add_heading("1. Company and mandate", level=1)
    doc.add_paragraph(C.COMPANY_CONTEXT)
    add_table(doc, ["Field", "Company detail"], [
        ("Department", a["case"]["department"]), ("Executive sponsor", a["case"]["sponsor"]),
        ("Business challenge", a["case"]["challenge"]), ("Planning horizon", "FY2027 with Q1 execution milestones"),
        ("Status", "Draft for classroom management review"),
    ], [1.6, 5.2])
    doc.add_heading("2. Evidence available", level=1)
    for source in a["case"]["sources"]:
        doc.add_paragraph(source, style="List Bullet")
    if a["num"] == 4:
        doc.add_paragraph("Roles requested by the teams for FY2027:")
        add_table(doc, ["Team", "Role", "How many", "Salary band", "Wanted by", "Hiring manager", "Why it is needed"],
                  LAB04_ROLES, [0.8, 1.3, 0.6, 1.1, 0.8, 1.1, 1.8])
        p = doc.add_paragraph(); r = p.add_run(LAB04_BUDGET); r.bold = True
    doc.add_heading("3. Required management outputs", level=1)
    add_table(doc, ["Output", "Acceptance evidence", "Owner"], [
        (value, "Source-linked, template-compliant and reviewed", a["case"]["sponsor"]) for value in a["case"]["outputs"]
    ], [2.2, 3.3, 1.3])
    doc.add_heading("4. Measures", level=1)
    add_table(doc, ["Measure", "Definition", "Source owner", "Review frequency"], [
        (metric, f"Company definition for {metric.lower()}", a["case"]["department"], "Monthly") for metric in a["case"]["metrics"]
    ], [1.6, 2.7, 1.5, 1.0])
    doc.add_heading("5. Controls and approval", level=1)
    for control in a["case"]["controls"]:
        doc.add_paragraph(control, style="List Bullet")
    callout(doc, "Evidence rule", "Facts cite their source; calculations cite workbook cells; assumptions name an owner and review date; final release requires the named sponsor.")
    doc.add_heading("6. Management questions", level=1)
    questions = [
        "What decision must be made, and by whom?", "Which evidence is authoritative?",
        "Which assumption could change the recommendation?", "What action, owner and date follow from approval?",
    ]
    for q in questions: doc.add_paragraph(q, style="List Number")
    path = os.path.join(folder, f"{stem(a)}-HR-Brief.docx")
    doc.save(path); print("Saved", path)


def sample_sections(a):
    by_lab = {
        4: [("Executive decision", "Prioritise repeat customers and profitable digital growth while pausing broad discount-led acquisition."),
            ("Marketing choices", "Focus owned channels, targeted lifecycle campaigns and higher-margin product stories."),
            ("90-day action", "Launch two measured pilots with stop/go rules, accountable owners and a reconciled budget.")],
        5: [("Strategic ambition", "Build a profitable omnichannel growth engine with a differentiated home-lifestyle proposition."),
            ("What we will do", "Fill the three roles the brief names, hold back the fourth, and confirm each start date with the hiring manager."),
            ("Governance", "Six initiatives report through a quarterly scorecard with explicit dependencies and risk owners.")],
        6: [("Reporting boundary", "Singapore retail, office and directly controlled e-commerce operations for FY2026."),
            ("Sustainability performance", "Energy and waste movements are reported with method notes and limitations; no unsupported green claim is made."),
            ("Flexible-work policy", "Eligibility is role-based, decisions are documented and exceptions have a transparent review route.")],
        7: [("CFO headline", "Revenue finished ahead of budget, but contribution quality differs materially by channel."),
            ("Driver", "Digital growth supported the top line while discount and fulfilment economics weakened selected contribution rates."),
            ("Decision", "Protect margin through targeted discount controls and channel-specific action owners.")],
        8: [("Executive story", "Approve three strategic choices, a focused marketing allocation and Q1 financial guardrails."),
            ("Evidence", "The storyline reconciles the approved strategy, marketing plan and financial dashboard."),
            ("Presentation standard", "Native charts, company layouts, conclusion-led titles and source notes keep the deck editable and defensible.")],
        9: [("Thread summary", "The Executive Committee needs the reviewed pack, three owned actions and a confirmed review date."),
            ("Draft response", "A concise reply confirms decisions and attachments and remains unsent until recipient and version checks are complete."),
            ("Meeting brief", "Purpose, agenda, pre-read and one decision required are explicit.")],
        10: [("Cowork hand-off", "The scoped project reconciles source files before generating a two-page management brief."),
             ("Discrepancy control", "Figures, owners, dates, versions and approvals are resolved in the authoritative source."),
             ("Native review", "Generated files return to Word, Excel and PowerPoint for tracked human review.")],
        11: [("Daily headline", "One margin exception and two overdue actions require management attention today."),
             ("Evidence", "KPI exceptions cite workbook cells and decisions cite approved Outlook messages."),
             ("Automation control", "The run is idempotent, creates a backup, logs changes and does not send email.")],
    }
    return by_lab.get(a["num"], [
        ("Management outcome", a["objective"]), ("Evidence boundary", "; ".join(a["case"]["sources"])),
        ("Approval", f"Final approval: {a['case']['sponsor']}"),
    ])


def save_work_sample(a, folder):
    doc = Document(); set_doc_defaults(doc); add_doc_brand(doc, f"{C.COURSE_CODE} · Claude-generated work sample")
    doc.add_paragraph(a["title"], style="Title")
    p = doc.add_paragraph(f"{C.COMPANY} · Illustrative reviewed output")
    p.runs[0].bold = True; p.runs[0].font.color.rgb = DocRGB.from_string(TEAL)
    callout(doc, "Management outcome", a["case"]["decision"], MINT)
    for title, body in sample_sections(a):
        doc.add_heading(title, level=1)
        doc.add_paragraph(body)
        if title.lower() in ("marketing choices", "strategic choices", "driver", "evidence"):
            add_table(doc, ["Finding / choice", "Evidence", "Owner", "Next review"], [
                (f"{title} item {i}", a["case"]["sources"][(i-1) % len(a["case"]["sources"])], a["case"]["department"], f"Q{i} FY2027") for i in range(1, 4)
            ], [2.0, 2.4, 1.5, 0.9])
    doc.add_heading("Source and approval note", level=1)
    doc.add_paragraph("Illustrative output generated for training from fictional Lumina Living sources. Figures and claims require source reconciliation. Release status: Pending named human approval.")
    path = os.path.join(folder, f"{stem(a)}-Claude-Generated-Work-Sample.docx")
    doc.save(path); print("Saved", path)


def save_prompt_template(a, folder):
    if a["num"] == 1:
        doc = Document(); set_doc_defaults(doc); add_doc_brand(doc, f"{C.COURSE_CODE} · Trainer guide")
        doc.add_paragraph("Lab 01 Trainer Demonstration Guide", style="Title")
        p = doc.add_paragraph("Teach three routes, practise one email, complete one checklist")
        p.runs[0].bold = True; p.runs[0].font.color.rgb = DocRGB.from_string(TEAL)
        callout(doc, "Teaching outcome", a["objective"], MINT)
        doc.add_heading("20-minute run sheet", level=1)
        add_table(doc, ["Time", "Trainer action", "Learner evidence"], [
            (timing, f"{action}: {teaching}", evidence)
            for timing, action, teaching, evidence in a["trainer_plan"]
        ], [1.0, 4.2, 2.0])
        doc.add_heading("Before class", level=1)
        for item in a["trainer_preclass"]: doc.add_paragraph(item, style="List Bullet")
        doc.add_heading("Keep out of Lab 01", level=1)
        for item in a["trainer_exclusions"]: doc.add_paragraph(item, style="List Bullet")
        doc.add_heading("Completion standard", level=1)
        doc.add_paragraph(a["test"])
        path = os.path.join(folder, "templates", "Lab-01-Trainer-Demonstration-Guide.docx")
        doc.save(path); print("Saved", path)
        return
    doc = Document(); set_doc_defaults(doc); add_doc_brand(doc, f"{C.COURSE_CODE} · Reusable template")
    doc.add_paragraph("Prompt and Review Contract", style="Title")
    p = doc.add_paragraph(f"Lab {a['num']} · {a['title']}"); p.runs[0].bold = True; p.runs[0].font.color.rgb = DocRGB.from_string(TEAL)
    for heading, hint in [
        ("1. Business result", "Decision, audience and artifact required"),
        ("2. Approved evidence", "Open files, tables, ranges, sections or messages Claude may use"),
        ("3. Constraints", "Scope, style, formulas, length, preservation rules and prohibited actions"),
        ("4. Output contract", "Required structure, format, citations and exception handling"),
        ("5. Verification", "Checks the user will perform independently"),
        ("6. Approval gate", "Named person who accepts, saves, writes, sends or releases"),
    ]:
        doc.add_heading(heading, level=1)
        callout(doc, "Complete this field", hint, LIGHT)
        doc.add_paragraph("\n")
    doc.add_heading("Review checklist", level=1)
    add_table(doc, ["Check", "Evidence", "Result", "Reviewer"], [
        ("Sources cited", "", "Pending", ""), ("Figures reconciled", "", "Pending", ""),
        ("Template preserved", "", "Pending", ""), ("Risks and assumptions labelled", "", "Pending", ""),
        ("Approval recorded", "", "Pending", ""),
    ], [2.0, 2.5, 1.1, 1.2])
    path = os.path.join(folder, "templates", "Prompt-and-Review-Template.docx")
    doc.save(path); print("Saved", path)


def style_header(ws, row=1):
    for cell in ws[row]:
        if cell.value is None: continue
        cell.font = Font(name="Arial", bold=True, color=WHITE)
        cell.fill = PatternFill("solid", fgColor=BLUE)
        cell.alignment = Alignment(wrap_text=True, vertical="center")


def style_sheet(ws, widths=None, freeze="A2"):
    ws.sheet_view.showGridLines = False
    if freeze: ws.freeze_panes = freeze
    for col, width in (widths or {}).items(): ws.column_dimensions[col].width = width
    ws.page_setup.orientation = "landscape"; ws.page_setup.fitToWidth = 1; ws.page_setup.fitToHeight = 0
    ws.sheet_properties.pageSetUpPr.fitToPage = True


MONTHS = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]


def add_readme_sheet(wb, a):
    ws = wb.active; ws.title = "Read_Me"; ws.sheet_view.showGridLines = False
    ws.column_dimensions["A"].width = 3; ws.column_dimensions["B"].width = 25; ws.column_dimensions["C"].width = 90
    ws.merge_cells("B2:C3"); ws["B2"] = f"{C.COMPANY.upper()} · LAB {a['num']:02d}"
    ws["B2"].font = Font(name="Arial", size=22, bold=True, color=WHITE); ws["B2"].fill = PatternFill("solid", fgColor=NAVY); ws["B2"].alignment = Alignment(vertical="center")
    rows = [("Activity", a["title"]), ("Decision", a["case"]["decision"]), ("Department", a["case"]["department"]),
            ("Sponsor", a["case"]["sponsor"]), ("Fictional data", "All names, messages and figures are synthetic training material."),
            ("Acceptance", a["test"])]
    for r, (label, value) in enumerate(rows, 5):
        ws[f"B{r}"] = label; ws[f"C{r}"] = value
        ws[f"B{r}"].font = Font(name="Arial", bold=True, color=BLUE); ws[f"B{r}"].fill = PatternFill("solid", fgColor=ICE)
        ws[f"C{r}"].alignment = Alignment(wrap_text=True, vertical="top")
    ws.print_area = "B2:C10"; ws.page_setup.fitToHeight = 1


def add_control_data(wb, a):
    ws = wb.create_sheet("Management_Control")
    style_sheet(ws, {"A": 12, "B": 24, "C": 16, "D": 16, "E": 16, "F": 16, "G": 18, "H": 18})
    headers = ["Month", "Measure", "Budget", "Actual", "Variance", "Variance_Pct", "Owner", "Status"]
    ws.append(headers); style_header(ws)
    random.seed(1970 + a["num"])
    row = 2
    for month_i, month in enumerate(MONTHS, 1):
        for metric_i, metric in enumerate(a["case"]["metrics"][:4], 1):
            budget = 80000 + a["num"] * 1700 + metric_i * 6200 + month_i * 1250
            factor = 0.91 + random.random() * 0.20
            actual = round(budget * factor, 2)
            ws.append([month, metric, budget, actual, f"=D{row}-C{row}", f'=IF(C{row}=0,0,E{row}/C{row})', a["case"]["department"], f'=IF(ABS(F{row})<=0.05,"On track",IF(F{row}>0.05,"Above","Attention"))'])
            ws[f"C{row}"].number_format = "$#,##0"; ws[f"D{row}"].number_format = "$#,##0"; ws[f"E{row}"].number_format = "$#,##0"; ws[f"F{row}"].number_format = "0.0%"
            row += 1
    table = Table(displayName=f"tblLab{a['num']:02d}Control", ref=f"A1:H{row-1}")
    table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True, showColumnStripes=False)
    ws.add_table(table); ws.auto_filter.ref = f"A1:H{row-1}"
    return ws


LAB02_USES = {
    "Staff records folder": "Employee records and contracts",
    "Hiring files": "CVs, interview notes and offers",
    "Payroll summary": "Salary and payroll data",
    "HR mailbox": "Questions from staff",
}


def add_source_register(wb, a):
    ws = wb.create_sheet("Where_Info_Is_Kept")
    style_sheet(ws, {"A": 7, "B": 34, "C": 24, "D": 18, "E": 18, "F": 18})
    ws.append(["Ref", "Where staff information is kept", "Who owns it", "What it is used for", "Last checked", "Status"]); style_header(ws)
    for i, source in enumerate(a["case"]["sources"], 1):
        ws.append([f"S{i:02d}", source, a["case"]["department"], LAB02_USES.get(source, "Staff information"), date(2026, 8, 12), "Approved for training"])
        ws[f"E{i+1}"].number_format = "dd-mmm-yyyy"
    return ws


# Lab 09 works on a local inbox instead of a real mailbox, so the lab runs on
# any computer with no Outlook, no work account and nothing ever sent.
LAB09_INBOX = [
    ("LL-MSG-001", "2026-08-10 08:41", "Priya Raman, Finance", "Carry-over leave — how many days?",
     "I have 7 days of annual leave left. The handbook says up to five days may be carried over but my manager says three. Which is right?"),
    ("LL-MSG-002", "2026-08-10 09:05", "Marcus Tan, Warehouse", "Probation extended without explanation",
     "My probation was due to end last Friday. My supervisor says it has been extended by a month but could not tell me why or who approved it."),
    ("LL-MSG-003", "2026-08-10 09:20", "Facilities", "Office lift maintenance this Saturday",
     "The lifts will be out of service on Saturday from 7am to 1pm. No action needed."),
    ("LL-MSG-004", "2026-08-10 10:02", "Siti Rahman, Online", "Request to work from home two days a week",
     "I would like to work from home on Tuesdays and Thursdays. I understand a flexible working policy is being drafted. Can I apply now or should I wait?"),
    ("LL-MSG-005", "2026-08-10 10:47", "Legal", "Sustainability wording needs review",
     "The draft sustainability section states an emissions reduction we cannot yet evidence. Please remove or qualify it before the pack is circulated."),
    ("LL-MSG-006", "2026-08-10 11:15", "IT Service Desk", "Scheduled password policy change",
     "Password rotation moves to 180 days from next month. This is for information only."),
    ("LL-MSG-007", "2026-08-10 11:52", "David Lim, Retail", "Notice period on resignation",
     "I am considering resigning. The handbook says one month after probation. Does that change if I am asked to leave earlier, and is pay in lieu possible?"),
    ("LL-MSG-008", "2026-08-10 13:30", "Agency partner", "Creative concepts for Q1",
     "Attaching three creative routes for the Q1 campaign. No response needed until the segments are confirmed."),
    ("LL-MSG-009", "2026-08-10 14:08", "Head of HR office", "Handbook sections needed by 4pm",
     "The Head of HR has asked for the probation and flexible working sections of the handbook by 4pm today."),
    ("LL-MSG-010", "2026-08-10 15:22", "Nurul Aziz, HR", "Policy draft circulation",
     "The HR policy draft is ready for review, but it should not be circulated until legal has confirmed the leave provisions."),
]


def add_inbox_sheet(wb):
    ws = wb.create_sheet("Staff_Messages", 1)
    style_sheet(ws, {"A": 14, "B": 16, "C": 24, "D": 40, "E": 70, "F": 22, "G": 40})
    ws.append(["Message_ID", "Received", "From", "Subject", "Message", "Category", "Draft_Reply"])
    style_header(ws)
    for row in LAB09_INBOX:
        ws.append(list(row) + ["", ""])
    for r in range(2, len(LAB09_INBOX) + 2):
        ws[f"E{r}"].alignment = Alignment(wrap_text=True, vertical="top")
        ws[f"G{r}"].alignment = Alignment(wrap_text=True, vertical="top")
        ws.row_dimensions[r].height = 42
    return ws


def add_permission_map(wb, a):
    """Lab 02's real deliverable: one row per source, learner completes the scope."""
    ws = wb.create_sheet("What_Claude_May_Do", 2)
    style_sheet(ws, {"A": 10, "B": 30, "C": 22, "D": 26, "E": 18, "F": 24, "G": 30})
    ws.append(["Ref", "Where staff information is kept", "Who uses it",
               "Read only or read and change", "Who owns it", "Check again by", "What to do if it is unavailable"])
    style_header(ws)
    for i, source in enumerate(a["case"]["sources"], 1):
        ws.append([f"S{i:02d}", source, "Business Transformation Office", "", "", "", ""])
    for r in range(2, len(a["case"]["sources"]) + 2):
        ws[f"B{r}"].alignment = Alignment(wrap_text=True, vertical="top")
        ws[f"G{r}"].alignment = Alignment(wrap_text=True, vertical="top")
        ws.row_dimensions[r].height = 30
    ws[f"A{len(a['case']['sources']) + 3}"] = "Complete the blank columns yourself. One row per source."
    return ws


def add_review_log(wb, title="Review_Log"):
    ws = wb.create_sheet(title)
    style_sheet(ws, {"A": 14, "B": 24, "C": 38, "D": 18, "E": 18, "F": 16})
    ws.append(["Date", "Check", "Evidence / change", "Owner", "Reviewer", "Status"]); style_header(ws)
    ws.append([date(2026, 8, 12), "Starter created", "Fictional training workbook generated", "Courseware team", "", "Open"])
    ws["A2"].number_format = "dd-mmm-yyyy"
    return ws


# Lab 01 shortlist. AI experience is written into the notes the way it appears
# on a real application, so the learner cannot shortlist by scanning one column.
LAB01_CANDIDATES = [
    ("C-001", "Amirah Bakar", "S8712xxxx", "9123 4xxx", "amirah.b@example.com", "Diploma in Business", 4,
     "Payroll and leave administration. Built team dashboards in Excel.", "Excel, Workday", 3600),
    ("C-002", "Daniel Ong", "S8534xxxx", "9234 5xxx", "daniel.ong@example.com", "Degree in HR Management", 7,
     "Ran hiring for two retail regions. Used Claude to draft job adverts and screen CVs.", "Recruitment, Claude, Excel", 5200),
    ("C-003", "Wei Ling Chua", "S8823xxxx", "9345 6xxx", "weiling.c@example.com", "Diploma in Psychology", 6,
     "Agency background, high-volume warehouse hiring. No AI tools used.", "Sourcing, interviewing", 4800),
    ("C-004", "Rajesh Kumar", "S9011xxxx", "9456 7xxx", "rajesh.k@example.com", "Degree in Statistics", 3,
     "Headcount reporting and turnover analysis. Automated monthly reports with ChatGPT and Python.", "Excel, Python, ChatGPT", 4500),
    ("C-005", "Siti Nurhaliza", "S9302xxxx", "9567 8xxx", "siti.n@example.com", "Diploma in Business Admin", 2,
     "Onboarding coordination and staff records. Familiar with Workday.", "Workday, MS Office", 3200),
    ("C-006", "Marcus Lee", "S8109xxxx", "9678 9xxx", "marcus.lee@example.com", "Degree in Education", 9,
     "Designed induction programmes. Piloted an AI assistant for course content.", "L&D design, LMS", 6000),
    ("C-007", "Priya Menon", "S8745xxxx", "9789 0xxx", "priya.m@example.com", "Degree in Law", 5,
     "Employee relations and grievance handling. No AI experience stated.", "ER casework, policy", 5000),
    ("C-008", "Jonathan Tay", "S8256xxxx", "9890 1xxx", "jon.tay@example.com", "Degree in Information Systems", 8,
     "HRIS migration lead. Built Claude-based workflows for policy Q&A.", "HRIS, Claude, SQL", 6500),
    ("C-009", "Farah Ismail", "S9204xxxx", "9901 2xxx", "farah.i@example.com", "Diploma in Mass Comm", 3,
     "Campus hiring and interview scheduling. Uses LinkedIn Recruiter daily.", "Sourcing, scheduling", 3400),
    ("C-010", "Kenneth Goh", "S8631xxxx", "9012 3xxx", "kenneth.g@example.com", "Degree in Economics", 6,
     "Salary benchmarking and bonus modelling. Advanced Excel, no AI tools.", "Excel, benchmarking", 5500),
    ("C-011", "Aisha Rahman", "S8918xxxx", "9123 5xxx", "aisha.r@example.com", "Diploma in HR", 4,
     "Leave and benefits administration. Attended an AI-for-HR short course last year.", "Benefits, MS Office", 3800),
    ("C-012", "Benjamin Sim", "S7922xxxx", "9234 6xxx", "ben.sim@example.com", "Degree in Business", 10,
     "Built the hiring function at a logistics firm. Trialled Claude for interview scorecards.", "TA strategy, Claude", 7000),
    ("C-013", "Nur Syafiqah", "S9508xxxx", "9345 7xxx", "nur.s@example.com", "Diploma in Business", 1,
     "Records filing and first-line staff queries. Studying for a CIPD certificate.", "MS Office", 2600),
    ("C-014", "Vincent Ng", "S7815xxxx", "9456 8xxx", "vincent.ng@example.com", "Degree in Management", 11,
     "Full-cycle HR for a 300-person retailer. No AI tools in current role.", "Generalist HR, ER", 7200),
    ("C-015", "Divya Pillai", "S8827xxxx", "9567 9xxx", "divya.p@example.com", "Masters in Analytics", 7,
     "Turnover prediction models. Uses Claude and Python for analysis and reporting.", "Python, Claude, Power BI", 6800),
    ("C-016", "Hafiz Osman", "S9412xxxx", "9678 0xxx", "hafiz.o@example.com", "Diploma in Business", 2,
     "Interview logistics and candidate communications.", "Scheduling, MS Office", 2900),
    ("C-017", "Grace Tan", "S8703xxxx", "9789 1xxx", "grace.tan@example.com", "Degree in Sociology", 6,
     "Restructuring and workforce planning. Experimented with generative AI for policy drafts.", "Workforce planning", 5600),
    ("C-018", "Ahmad Zulkifli", "S8940xxxx", "9890 2xxx", "ahmad.z@example.com", "Diploma in Accounting", 5,
     "Monthly payroll for 400 staff. Strong Excel, no AI exposure.", "Payroll, Excel", 4200),
    ("C-019", "Michelle Koh", "S9119xxxx", "9901 3xxx", "michelle.k@example.com", "Degree in Communications", 3,
     "Training coordination and feedback analysis. Used AI tools to summarise course feedback.", "L&D admin, survey tools", 3600),
    ("C-020", "Suresh Nair", "S7736xxxx", "9012 4xxx", "suresh.n@example.com", "Degree in Industrial Relations", 12,
     "Union negotiation and ER casework. No AI tools used.", "IR, negotiation", 7500),
    ("C-021", "Lim Jia Hui", "S8852xxxx", "9123 6xxx", "jiahui.lim@example.com", "Degree in Business", 4,
     "Onboarding and probation tracking. Built an AI chatbot for FAQs in a previous role.", "Onboarding, chatbots", 4000),
    ("C-022", "Nadia Hassan", "S8347xxxx", "9234 7xxx", "nadia.h@example.com", "Degree in Psychology", 8,
     "Executive search background. No AI experience stated.", "Executive search", 6200),
    ("C-023", "Terrence Wong", "S9425xxxx", "9345 8xxx", "terrence.w@example.com", "Diploma in Data Analytics", 2,
     "Headcount and cost reporting. Uses Claude in Excel for formula help.", "Excel, Claude", 3300),
    ("C-024", "Yasmin Abdullah", "S8218xxxx", "9456 9xxx", "yasmin.a@example.com", "Degree in HR Management", 9,
     "Performance management and succession planning. No AI tools in current role.", "Performance, succession", 6400),
]



# Lab 02 works on a real staff list so learners can chart and analyse it.
LAB02_STAFF = [
    ("Retail", "Store Assistant", 18, 2450, 4, 12),
    ("Retail", "Store Supervisor", 6, 3800, 1, 9),
    ("Retail", "Visual Merchandiser", 3, 3400, 0, 6),
    ("Online", "Customer Service", 12, 2900, 3, 15),
    ("Online", "Content Executive", 5, 3600, 1, 7),
    ("Online", "Digital Analyst", 4, 4800, 0, 3),
    ("Warehouse", "Picker Packer", 22, 2300, 6, 21),
    ("Warehouse", "Forklift Operator", 7, 2900, 2, 10),
    ("Warehouse", "Shift Supervisor", 5, 3900, 1, 8),
    ("Office", "Finance Executive", 6, 4200, 0, 4),
    ("Office", "HR Executive", 4, 3900, 1, 5),
    ("Office", "IT Support", 3, 4400, 0, 2),
]


def add_staff_list(wb):
    ws = wb.create_sheet("Staff_List", 1)
    style_sheet(ws, {"A": 8, "B": 16, "C": 24, "D": 12, "E": 16, "F": 16, "G": 18})
    ws.append(["Ref", "Team", "Role", "Headcount",
               "Average monthly salary", "Left in last 12 months", "Open days to fill"])
    style_header(ws)
    for i, (team, role, head, salary, leavers, days) in enumerate(LAB02_STAFF, 1):
        ws.append([f"R{i:02d}", team, role, head, salary, leavers, days])
    for r in range(2, len(LAB02_STAFF) + 2):
        ws[f"E{r}"].number_format = "#,##0"
    ws.freeze_panes = "A2"
    return ws



# Lab 11 (Skills) needs a fresh applicant list so the saved Skill can be run on
# data the learner has not already sorted by hand.
LAB11_APPLICANTS = [
    ("A-01", "Rachel Sim", "HR Executive", 5, "Leave administration and staff records. Uses Claude weekly to draft policy replies."),
    ("A-02", "Imran Yusof", "Recruitment Lead", 8, "Volume hiring for logistics. No AI tools used."),
    ("A-03", "Chloe Tan", "HR Analyst", 3, "Headcount reporting. Built turnover dashboards with ChatGPT assistance."),
    ("A-04", "Ganesh Raj", "L&D Executive", 4, "Course scheduling and feedback. Completed an AI fundamentals certificate."),
    ("A-05", "Melissa Chong", "HR Business Partner", 7, "Employee relations casework. No AI experience stated."),
    ("A-06", "Zainal Abidin", "Payroll Specialist", 6, "Payroll for 500 staff. Trialled Copilot for spreadsheet formulas."),
    ("A-07", "Vanessa Lim", "Talent Acquisition", 4, "Campus and graduate hiring. Uses AI screening tools daily."),
    ("A-08", "Haziq Rahman", "HR Assistant", 1, "Filing and first-line queries. Studying HR analytics part-time."),
    ("A-09", "Serena Wong", "People Analytics", 6, "Workforce planning models. Uses Claude and Python for reporting."),
    ("A-10", "Kumar Selvam", "HR Manager", 10, "Full-cycle HR at a manufacturer. No AI tools in current role."),
    ("A-11", "Adeline Foo", "Recruitment Executive", 2, "Interview scheduling and candidate care."),
    ("A-12", "Faizal Omar", "HR Systems Analyst", 9, "HRIS administration. Automated onboarding checks with an AI assistant."),
    ("A-13", "Jasmine Neo", "Compensation Analyst", 5, "Benchmarking and bonus modelling. Attended an AI-for-HR webinar."),
    ("A-14", "Ridwan Salleh", "L&D Manager", 11, "Leadership programmes. No AI tools used."),
    ("A-15", "Cheryl Ang", "HR Executive", 3, "Onboarding and probation. Uses Claude in Word to draft letters."),
    ("A-16", "Nurul Huda", "HR Business Partner", 6, "Restructuring support. Experimented with AI for policy summaries."),
    ("A-17", "Desmond Yeo", "Recruitment Coordinator", 2, "Job posting and applicant tracking. No AI experience stated."),
    ("A-18", "Priyanka Devi", "HR Analyst", 4, "Attrition analysis. Reads AI research but has not used the tools at work."),
]


def build_skill_applicants(a):
    wb = Workbook(); add_readme_sheet(wb, a)
    ws = wb.create_sheet("Applicants")
    style_sheet(ws, {"A": 8, "B": 20, "C": 24, "D": 9, "E": 66, "F": 14, "G": 30})
    ws.append(["Ref", "Name", "NRIC", "Contact", "Email", "Highest qualification",
               "Years", "Experience notes", "Skills", "Expected salary",
               "Interview", "Why"])
    style_header(ws)
    for row in LAB11_APPLICANTS:
        ws.append(list(row) + ["", ""])
    for r in range(2, len(LAB11_APPLICANTS) + 2):
        ws[f"E{r}"].alignment = Alignment(wrap_text=True, vertical="top")
        ws[f"G{r}"].alignment = Alignment(wrap_text=True, vertical="top")
        ws.row_dimensions[r].height = 30
    ws.freeze_panes = "A2"
    return wb


def build_surface_readiness_workbook(a):
    """Lab 01 works on a real shortlist so learners see a useful answer at once."""
    wb = Workbook(); add_readme_sheet(wb, a)
    readme = wb["Read_Me"]
    readme["B12"] = "Complete"
    readme["C12"] = ("Shortlist the candidates with hands-on Claude or AI experience, "
                     "then compare the answer you get in Excel with the answer you get "
                     "in the Claude Desktop app.")
    readme["B12"].font = Font(name="Arial", bold=True, color=BLUE)
    readme["B12"].fill = PatternFill("solid", fgColor=ICE)
    readme["C12"].alignment = Alignment(wrap_text=True)
    readme.print_area = "B2:C12"
    readme.page_setup.orientation = "portrait"
    readme.page_setup.fitToWidth = 1; readme.page_setup.fitToHeight = 1
    readme.sheet_properties.pageSetUpPr.fitToPage = True

    ws = wb.create_sheet("Candidates")
    style_sheet(ws, {"A": 8, "B": 18, "C": 12, "D": 12, "E": 24, "F": 24,
                     "G": 7, "H": 54, "I": 24, "J": 14, "K": 11, "L": 30})
    ws.page_setup.orientation = "landscape"
    ws.append(["Ref", "Name", "NRIC", "Contact", "Email", "Highest qualification",
               "Years", "Experience notes", "Skills", "Expected salary",
               "Interview", "Why"])
    style_header(ws)
    for row in LAB01_CANDIDATES:
        ws.append(list(row) + ["", ""])
    for r in range(2, len(LAB01_CANDIDATES) + 2):
        for col in ("H", "I", "L"):
            ws[f"{col}{r}"].alignment = Alignment(wrap_text=True, vertical="top")
        ws[f"J{r}"].number_format = "#,##0"
        ws.row_dimensions[r].height = 34
    ws.freeze_panes = "A2"
    return wb

def build_finance_workbook(a):
    wb = Workbook(); add_readme_sheet(wb, a)
    dash = wb.create_sheet("Dashboard", 0); dash.sheet_view.showGridLines = False
    for col in range(1, 15): dash.column_dimensions[chr(64 + col)].width = 10
    dash.merge_cells("B2:M3"); dash["B2"] = "LUMINA LIVING · FY2026 FINANCIAL DASHBOARD"
    dash["B2"].font = Font(name="Arial", size=21, bold=True, color=WHITE); dash["B2"].fill = PatternFill("solid", fgColor=NAVY); dash["B2"].alignment = Alignment(horizontal="center", vertical="center")
    tx = wb.create_sheet("Transactions")
    style_sheet(tx, {"A": 14, "B": 13, "C": 11, "D": 13, "E": 14, "F": 22, "G": 12, "H": 10, "I": 13, "J": 13, "K": 14, "L": 13, "M": 14, "N": 15})
    headers = ["Txn_ID", "Date", "Month", "Region", "Channel", "Product", "Category", "Units", "Unit_Price", "Discount_Pct", "Revenue", "Unit_Cost", "Gross_Profit", "Fulfilment_Cost"]
    tx.append(headers); style_header(tx)
    products = {"Aurora Desk Lamp": ("Lighting", 89, 42), "Halo Floor Lamp": ("Lighting", 159, 79), "Luna Throw": ("Living", 69, 31), "Nest Basket": ("Living", 55, 24), "Arc Side Table": ("Living", 139, 72), "Focus Monitor Stand": ("Workspace", 79, 36)}
    regions = ["North", "South", "East", "West"]; channels = ["Retail", "Online", "Marketplace"]
    random.seed(197); start = date(2026, 1, 1)
    for i in range(1, 721):
        dt = start + timedelta(days=random.randint(0, 364)); month = dt.strftime("%b")
        region = random.choices(regions, [31, 23, 27, 19])[0]; channel = random.choices(channels, [42, 37, 21])[0]
        product = random.choice(list(products)); category, price, cost = products[product]
        units = random.randint(1, 8); discount = random.choices([0, .05, .10, .15, .20], [20, 22, 31, 20, 7])[0]
        r = i + 1; fulfil = 6 if channel == "Retail" else (11 if channel == "Online" else 15)
        tx.append([f"LL-{i:05d}", dt, month, region, channel, product, category, units, price, discount, f"=H{r}*I{r}*(1-J{r})", cost, f"=K{r}-(H{r}*L{r})", f"=H{r}*{fulfil}"])
        for c in (9, 11, 12, 13, 14): tx.cell(r, c).number_format = "$#,##0.00"
        tx.cell(r, 10).number_format = "0%"; tx.cell(r, 2).number_format = "dd-mmm-yyyy"
    table = Table(displayName="tblFinance", ref=f"A1:N{tx.max_row}"); table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True, showColumnStripes=False); tx.add_table(table)
    budget = wb.create_sheet("Budget"); style_sheet(budget, {"A": 12, "B": 18, "C": 18, "D": 18, "E": 18})
    budget.append(["Month", "Revenue_Budget", "Gross_Profit_Budget", "Marketing_Budget", "Operating_Cost_Budget"]); style_header(budget)
    for i, month in enumerate(MONTHS, 2):
        rev = 72000 + (i - 2) * 2600
        budget.append([month, rev, rev * .43, 9000 + (i % 3) * 1200, 21000 + (i % 4) * 700])
        for c in range(2, 6): budget.cell(i, c).number_format = "$#,##0"
    assumptions = wb.create_sheet("Assumptions"); style_sheet(assumptions, {"A": 26, "B": 16, "C": 48})
    assumptions.append(["Assumption", "Base value", "Definition / control"]); style_header(assumptions)
    for row in [("Unit growth", .08, "Applied only to FY2027 scenario"), ("Average discount", .09, "Weighted average target"), ("Unit-cost inflation", .04, "Supplier planning assumption"), ("Marketing spend change", .06, "Versus FY2026 budget")]: assumptions.append(row)
    for r in range(2, 6): assumptions[f"B{r}"].number_format = "0.0%"
    ana = wb.create_sheet("Analysis"); ana.sheet_view.showGridLines = False
    for col, width in {"A": 3, "B": 22, "C": 18, "D": 18, "E": 4, "F": 17, "G": 17, "H": 17, "I": 17}.items(): ana.column_dimensions[col].width = width
    ana.merge_cells("B2:I2"); ana["B2"] = "FY2026 ACTUAL VS BUDGET ANALYSIS"; ana["B2"].font = Font(name="Arial", size=20, bold=True, color=WHITE); ana["B2"].fill = PatternFill("solid", fgColor=NAVY)
    kpis = [("Revenue", "=SUM(Transactions!K2:K721)"), ("Gross Profit", "=SUM(Transactions!M2:M721)"), ("Gross Margin", "=IF(C4=0,0,C5/C4)"), ("Operating Contribution", "=C5-SUM(Transactions!N2:N721)-SUM(Budget!D2:D13)")]
    for r, (label, formula) in enumerate(kpis, 4):
        ana[f"B{r}"] = label; ana[f"C{r}"] = formula; ana[f"B{r}"].font = Font(name="Arial", bold=True, color=BLUE); ana[f"B{r}"].fill = PatternFill("solid", fgColor=ICE); ana[f"C{r}"].number_format = "0.0%" if label == "Gross Margin" else "$#,##0"
    ana.append([])
    for cell, value in [("B10", "Month"), ("C10", "Actual Revenue"), ("D10", "Budget Revenue"), ("F10", "Channel"), ("G10", "Revenue"), ("H10", "Gross Profit"), ("I10", "Contribution")]: ana[cell] = value
    for cell in ["B10", "C10", "D10", "F10", "G10", "H10", "I10"]: ana[cell].font = Font(name="Arial", bold=True, color=WHITE); ana[cell].fill = PatternFill("solid", fgColor=TEAL)
    for idx, month in enumerate(MONTHS, 11):
        ana[f"B{idx}"] = month; ana[f"C{idx}"] = f'=SUMIF(Transactions!$C$2:$C$721,B{idx},Transactions!$K$2:$K$721)'; ana[f"D{idx}"] = f'=INDEX(Budget!$B$2:$B$13,MATCH(B{idx},Budget!$A$2:$A$13,0))'
        ana[f"C{idx}"].number_format = "$#,##0"; ana[f"D{idx}"].number_format = "$#,##0"
    for idx, channel in enumerate(channels, 11):
        ana[f"F{idx}"] = channel; ana[f"G{idx}"] = f'=SUMIF(Transactions!$E$2:$E$721,F{idx},Transactions!$K$2:$K$721)'; ana[f"H{idx}"] = f'=SUMIF(Transactions!$E$2:$E$721,F{idx},Transactions!$M$2:$M$721)'; ana[f"I{idx}"] = f'=H{idx}-SUMIF(Transactions!$E$2:$E$721,F{idx},Transactions!$N$2:$N$721)'
        for c in "GHI": ana[f"{c}{idx}"].number_format = "$#,##0"
    # Dashboard KPI cards
    cards = [("B5:D5", "B6:D7", "REVENUE", "=Analysis!C4", "$#,##0"), ("E5:G5", "E6:G7", "GROSS PROFIT", "=Analysis!C5", "$#,##0"), ("H5:J5", "H6:J7", "GROSS MARGIN", "=Analysis!C6", "0.0%"), ("K5:M5", "K6:M7", "OPERATING CONTRIBUTION", "=Analysis!C7", "$#,##0")]
    for label_rng, value_rng, label, formula, fmt in cards:
        dash.merge_cells(label_rng); dash.merge_cells(value_rng)
        l = dash[label_rng.split(":")[0]]; v = dash[value_rng.split(":")[0]]
        l.value = label; v.value = formula
        for cell in (l, v): cell.fill = PatternFill("solid", fgColor=ICE); cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        l.font = Font(name="Arial", size=10, bold=True, color=BLUE); v.font = Font(name="Arial", size=18, bold=True, color=NAVY); v.number_format = fmt
    line = LineChart(); line.title = "Monthly revenue: actual vs budget"; line.style = 13; line.y_axis.title = "Revenue ($)"; line.x_axis.title = "Month"
    line.add_data(Reference(ana, min_col=3, max_col=4, min_row=10, max_row=22), titles_from_data=True); line.set_categories(Reference(ana, min_col=2, min_row=11, max_row=22)); line.width = 11.5; line.height = 7.2; dash.add_chart(line, "B10")
    bar = BarChart(); bar.type = "bar"; bar.title = "Operating contribution by channel"; bar.style = 10; bar.x_axis.title = "Contribution ($)"
    bar.add_data(Reference(ana, min_col=9, min_row=10, max_row=13), titles_from_data=True); bar.set_categories(Reference(ana, min_col=6, min_row=11, max_row=13)); bar.width = 11.5; bar.height = 7.2; dash.add_chart(bar, "H10")
    var = BarChart(); var.type = "col"; var.title = "Monthly budget variance"; var.style = 11
    # Use two series for transparency; management can see actual and budget side-by-side.
    var.add_data(Reference(ana, min_col=3, max_col=4, min_row=10, max_row=22), titles_from_data=True); var.set_categories(Reference(ana, min_col=2, min_row=11, max_row=22)); var.width = 12.2; var.height = 7.0; dash.add_chart(var, "B25")
    dash["H25"] = "DEFINITIONS & CONTROLS"; dash["H25"].font = Font(name="Arial", bold=True, color=WHITE); dash["H25"].fill = PatternFill("solid", fgColor=TEAL)
    for i, text in enumerate(["Gross margin = Gross profit / Revenue", "Operating contribution deducts fulfilment and marketing budget", "Scenario assumptions remain separate from actuals", "All data is fictional training material", "Last refreshed: 12 Aug 2026"], 26): dash[f"H{i}"] = text
    dash.print_area = "B2:M34"; dash.page_setup.orientation = "landscape"; dash.page_setup.fitToWidth = 1; dash.page_setup.fitToHeight = 1; dash.sheet_properties.pageSetUpPr.fitToPage = True
    add_source_register(wb, a); add_review_log(wb, "Audit_Log")
    wb.calculation.fullCalcOnLoad = True; wb.calculation.forceFullCalc = True; wb.calculation.calcMode = "auto"
    return wb


def save_workbook(a, folder):
    if a["num"] == 11:
        wb = build_team_quarter(a, "Q1")
    elif a["num"] == 12:
        wb = build_team_quarter(a, "Q2")
    elif a["num"] == 15:
        wb = build_q1_staff(a)
    elif a["num"] == 1:
        wb = build_surface_readiness_workbook(a)
    elif a["num"] == 11:
        wb = build_skill_applicants(a)
    elif a["num"] == 6:
        wb = build_people_numbers(a)
    else:
        wb = Workbook(); add_readme_sheet(wb, a); add_control_data(wb, a); add_source_register(wb, a); add_review_log(wb)
        if a["num"] == 2: add_staff_list(wb); add_permission_map(wb, a)
        if a["num"] == 8: add_inbox_sheet(wb)
        # Put a simple management chart on a Summary sheet for realistic use.
        summary = wb.create_sheet("Summary", 0); summary.sheet_view.showGridLines = False
        summary.merge_cells("B2:J3"); summary["B2"] = f"{C.COMPANY.upper()} · {a['title'].upper()}"
        summary["B2"].font = Font(name="Arial", size=18, bold=True, color=WHITE); summary["B2"].fill = PatternFill("solid", fgColor=NAVY); summary["B2"].alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        summary["B5"] = "Decision"; summary["C5"] = a["case"]["decision"]; summary["B5"].font = Font(name="Arial", bold=True, color=BLUE); summary["B5"].fill = PatternFill("solid", fgColor=ICE); summary["C5"].alignment = Alignment(wrap_text=True)
        summary.column_dimensions["B"].width = 18; summary.column_dimensions["C"].width = 78
        summary["B8"] = "Use Management_Control for detailed fictional measures and Review_Log for approvals."
        summary.merge_cells("B8:J9"); summary["B8"].alignment = Alignment(wrap_text=True, vertical="center")
        summary.print_area = "B2:J12"; summary.page_setup.orientation = "landscape"; summary.page_setup.fitToWidth = 1; summary.page_setup.fitToHeight = 1; summary.sheet_properties.pageSetUpPr.fitToPage = True
    filename = (
        f"{stem(a)}-Q1-Staff.xlsx" if a["num"] == 15
        else f"{stem(a)}-Candidates.xlsx" if a["num"] == 1
        else f"{stem(a)}-Staff-Information.xlsx" if a["num"] == 2
        else f"{stem(a)}-Q1-Teams.xlsx" if a["num"] == 11
        else f"{stem(a)}-Q2-Teams.xlsx" if a["num"] == 12
        else f"{stem(a)}-Q1-Staff.xlsx" if a["num"] in (3, 4)
        else f"{stem(a)}-New-Applicants.xlsx" if a["num"] == 11
        else f"{stem(a)}-Candidates.xlsx" if a["num"] == 1
        else f"{stem(a)}-Staff-Information.xlsx" if a["num"] == 2
        else f"{stem(a)}-People-Numbers.xlsx" if a["num"] == 6
        else f"{stem(a)}-Staff-Questions.xlsx" if a["num"] == 8
        else f"{stem(a)}-Working-Workbook.xlsx"
    )
    path = os.path.join(folder, filename)
    wb.save(path); print("Saved", path)


def save_decision_log(a, folder):
    wb = Workbook(); ws = wb.active; ws.title = "Decision_Log"
    style_sheet(ws, {"A": 12, "B": 34, "C": 28, "D": 24, "E": 18, "F": 15, "G": 16})
    ws.append(["Decision_ID", "Decision / approval", "Evidence", "Options considered", "Owner", "Due date", "Status"]); style_header(ws)
    for i, output in enumerate(a["case"]["outputs"], 1):
        ws.append([f"D-{a['num']:02d}-{i:02d}", f"Approve {output}", "; ".join(a["case"]["sources"][:2]), "Approve / revise / defer", a["case"]["sponsor"], date(2026, 9, min(28, 5 + i * 4)), "Pending"])
        ws[f"F{i+1}"].number_format = "dd-mmm-yyyy"
    ws2 = wb.create_sheet("Approval_Checklist"); style_sheet(ws2, {"A": 30, "B": 20, "C": 24, "D": 18})
    ws2.append(["Control", "Evidence location", "Reviewer", "Status"]); style_header(ws2)
    for control in a["case"]["controls"]:
        ws2.append([control, "", "", "Pending"])
    path = os.path.join(folder, "templates", "Decision-and-Approval-Log.xlsx")
    wb.save(path); print("Saved", path)


def ppt_rect(slide, x, y, w, h, fill, radius=True, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, PIn(x), PIn(y), PIn(w), PIn(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor.from_string(fill); shape.line.color.rgb = RGBColor.from_string(line or fill)
    return shape


def ppt_text(slide, x, y, w, h, text, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PIn(x), PIn(y), PIn(w), PIn(h)); tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    box.text_frame.margin_left = 0; box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0; box.text_frame.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text; run.font.name = "Arial"; run.font.size = PPt(size); run.font.bold = bold; run.font.color.rgb = RGBColor.from_string(color)
    return box


def slide_base(prs, title, kicker, number):
    s = prs.slides.add_slide(prs.slide_layouts[6]); ppt_rect(s, 0, 0, 13.333, 7.5, WHITE, False)
    ppt_rect(s, 0, 0, .18, 7.5, TEAL, False); ppt_text(s, .75, .28, 11.6, .3, kicker.upper(), 11, TEAL, True)
    ppt_text(s, .75, .72, 11.8, .72, title, 27, NAVY, True); ppt_rect(s, .75, 1.62, 11.8, .02, LINE, False)
    ppt_text(s, .75, 7.08, 8.5, .2, f"{C.COMPANY} · Fictional training material", 8.5, GREY)
    ppt_text(s, 12.0, 7.05, .55, .22, str(number), 9, GREY, True, PP_ALIGN.RIGHT)
    return s


def add_cards(slide, items, y=2.0):
    cols = 2; card_w = 5.55; card_h = 1.55; gap_x = .55; gap_y = .38
    colors = [BLUE, TEAL, VIOLET, AMBER]
    for i, (head, body) in enumerate(items[:4]):
        row, col = divmod(i, cols); x = .85 + col * (card_w + gap_x); yy = y + row * (card_h + gap_y)
        ppt_rect(slide, x, yy, card_w, card_h, LIGHT, True, LINE); ppt_rect(slide, x, yy, .12, card_h, colors[i], False)
        ppt_text(slide, x + .28, yy + .13, card_w - .48, .38, head, 16, colors[i], True)
        ppt_text(slide, x + .28, yy + .56, card_w - .48, .78, body, 12.5, NAVY)


def add_flow(slide, steps, y=2.65):
    n = len(steps); gap = .22; total = 11.55; w = (total - gap * (n - 1)) / n
    for i, step in enumerate(steps):
        x = .88 + i * (w + gap); color = [BLUE, TEAL, VIOLET, AMBER, BLUE][i % 5]
        ppt_rect(slide, x, y, w, 1.6, LIGHT, True, LINE); ppt_text(slide, x + .12, y + .12, .42, .42, f"{i+1:02d}", 12, color, True)
        ppt_text(slide, x + .12, y + .58, w - .24, .7, step, 12.5, NAVY, True, PP_ALIGN.CENTER)
        if i < n - 1: ppt_text(slide, x + w, y + .55, gap, .45, "›", 22, GREY, True, PP_ALIGN.CENTER)


def add_native_chart(slide, x, y, w, h, title, categories, series):
    data = CategoryChartData(); data.categories = categories
    for name, values in series: data.add_series(name, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, PIn(x), PIn(y), PIn(w), PIn(h), data).chart
    chart.has_title = True; chart.chart_title.text_frame.text = title; chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    return chart


def save_presentation(a, folder):
    prs = Presentation(); prs.slide_width = PIn(13.333); prs.slide_height = PIn(7.5)
    cover = prs.slides.add_slide(prs.slide_layouts[6]); ppt_rect(cover, 0, 0, 13.333, 7.5, WHITE, False); ppt_rect(cover, 0, 0, .22, 7.5, TEAL, False)
    ppt_text(cover, .9, 1.0, 10.8, .35, f"{C.COMPANY.upper()} · FY2027", 14, TEAL, True)
    ppt_text(cover, .9, 1.65, 11.3, 1.45, a["title"], 36, NAVY, True)
    ppt_text(cover, .9, 3.45, 10.8, .8, a["case"]["decision"], 19, GREY)
    ppt_rect(cover, .9, 5.15, 3.7, .72, ICE, True); ppt_text(cover, 1.15, 5.28, 3.2, .4, f"LAB {a['num']:02d} · EDITABLE STARTER", 11, BLUE, True, PP_ALIGN.CENTER)
    ppt_text(cover, .9, 6.7, 10.8, .25, "Fictional company training material · Pending management approval", 10, GREY)
    s = slide_base(prs, "The business decision", "Executive context", 2)
    callouts = [("Challenge", a["case"]["challenge"]), ("Sponsor", a["case"]["sponsor"]), ("Outputs", "; ".join(a["case"]["outputs"])), ("Controls", "; ".join(a["case"]["controls"]))]
    add_cards(s, callouts)
    s = slide_base(prs, "Evidence-to-decision process", "Native process map", 3); add_flow(s, a["deck_flow"])
    s = slide_base(prs, "Management evidence", "Editable visual", 4)
    add_native_chart(s, .8, 1.95, 7.5, 4.75, "Illustrative monthly performance", MONTHS[:6], [("Budget", [82, 85, 88, 91, 94, 98]), ("Actual", [79, 87, 90, 88, 99, 104])])
    ppt_rect(s, 8.65, 1.95, 3.8, 4.75, LIGHT, True, LINE); ppt_text(s, 8.95, 2.25, 3.2, .35, "WHAT MANAGEMENT NEEDS", 11, TEAL, True)
    y = 2.85
    for metric in a["case"]["metrics"][:4]:
        ppt_text(s, 8.95, y, 3.05, .45, "• " + metric, 15, NAVY, True); y += .72
    s = slide_base(prs, "Professional practice", "Decision and control cards", 5); add_cards(s, a["deck_cards"])
    s = slide_base(prs, "Decision and approval", "Executive close", 6)
    cards = [("Decision required", a["case"]["decision"]), ("Evidence check", "Reconcile sources, cells, versions and owners."), ("Approval", f"Named sponsor: {a['case']['sponsor']}"), ("Next step", a["build"])]
    add_cards(s, cards)
    # Lab 8 carries a richer sample with an integrated strategic/marketing story.
    if a["num"] == 8:
        s = slide_base(prs, "Profitable growth requires three choices", "Strategy", 7)
        add_cards(s, [("Curated differentiation", "Invest in ranges where design, advice and margin reinforce one another."), ("Connected customer journey", "Use owned channels to move from discovery to repeat purchase."), ("Commercial discipline", "Tie campaign, discount and fulfilment decisions to contribution."), ("Explicit non-priority", "Avoid broad discount-led growth without a measurable stop/go rule.")])
        s = slide_base(prs, "Marketing allocation follows customer and margin evidence", "Marketing", 8)
        add_native_chart(s, .8, 1.95, 7.6, 4.75, "Illustrative channel contribution", ["Retail", "Online", "Marketplace"], [("Revenue", [420, 360, 220]), ("Contribution", [165, 132, 55])])
        ppt_rect(s, 8.75, 1.95, 3.7, 4.75, LIGHT, True, LINE); ppt_text(s, 9.0, 2.25, 3.15, .45, "DECISION LOGIC", 12, TEAL, True)
        ppt_text(s, 9.0, 3.0, 3.1, 2.9, "1  Protect profitable retail\n\n2  Scale owned digital\n\n3  Tighten marketplace economics", 17, NAVY, True)
        s = slide_base(prs, "Q1 roadmap converts choices into accountable work", "Execution", 9); add_flow(s, ["Confirm owners", "Launch two pilots", "Review economics", "Scale or stop", "Report to ExCo"])
        s = slide_base(prs, "Approve the plan with four guardrails", "Executive decision", 10); add_cards(s, [("Budget", "Release in stages against evidence."), ("Margin", "No growth target without contribution control."), ("Capability", "Fund analytics and lifecycle execution."), ("Governance", "Monthly review; quarterly choice reset.")])
    # Uniform Fade transition.
    from lxml import etree
    from pptx.oxml.ns import qn
    xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>'
    for slide in prs.slides:
        if slide.element.find(qn("p:transition")) is None: slide.element.append(etree.fromstring(xml))
    path = os.path.join(folder, f"{stem(a)}-Leadership-Update.pptx" if a["num"] == 7 else f"{stem(a)}-Executive-Starter.pptx")
    prs.save(path); print("Saved", path)


def save_lab11_automation(a, folder):
    auto = os.path.join(folder, "automation"); inputs = os.path.join(folder, "inputs"); outputs = os.path.join(folder, "outputs")
    os.makedirs(auto, exist_ok=True); os.makedirs(inputs, exist_ok=True); os.makedirs(outputs, exist_ok=True)
    updater = '''#!/usr/bin/env python3
"""Safe fictional daily-control updater. Review before use."""
import argparse, csv, os, shutil
from datetime import datetime
from openpyxl import load_workbook

p=argparse.ArgumentParser(); p.add_argument("--input",required=True); p.add_argument("--workbook",required=True); p.add_argument("--output",required=True); p.add_argument("--dry-run",action="store_true"); args=p.parse_args()
wb=load_workbook(args.workbook); required={"Management_Control","Review_Log"}; missing=required-set(wb.sheetnames)
if missing: raise SystemExit(f"Missing sheets: {sorted(missing)}")
rows=list(csv.DictReader(open(args.input,encoding="utf-8")))
ws=wb["Management_Control"]; keys={(ws.cell(r,1).value,ws.cell(r,2).value):r for r in range(2,ws.max_row+1)}
changes=[]
for item in rows:
    key=(item["Month"],item["Measure"]); row=keys.get(key)
    if not row: continue
    old=ws.cell(row,4).value; new=float(item["Actual"])
    if old!=new: changes.append((row,old,new))
if args.dry_run:
    print(f"DRY RUN: {len(changes)} proposed changes"); [print(x) for x in changes]; raise SystemExit(0)
backup=args.workbook+"."+datetime.now().strftime("%Y%m%d-%H%M%S")+".bak"; shutil.copy2(args.workbook,backup)
for row,old,new in changes: ws.cell(row,4).value=new
log=wb["Review_Log"]; log.append([datetime.now(),"Daily update",f"{len(changes)} actual values updated",os.getenv("USER","training-user"),"","Pending review"])
os.makedirs(os.path.dirname(args.output) or ".",exist_ok=True); wb.save(args.output); print(args.output)
'''
    brief = '''#!/usr/bin/env python3
"""Generate a source-linked fictional daily management brief."""
import argparse, json
from docx import Document
from openpyxl import load_workbook

p=argparse.ArgumentParser(); p.add_argument("--workbook",required=True); p.add_argument("--mail",required=True); p.add_argument("--template",required=True); p.add_argument("--output",required=True); args=p.parse_args()
wb=load_workbook(args.workbook,data_only=False); ws=wb["Management_Control"]
exceptions=[]
for r in range(2,ws.max_row+1):
    status=ws.cell(r,8).value
    if status and "Attention" in str(status): exceptions.append((ws.cell(r,1).value,ws.cell(r,2).value,f"Management_Control!H{r}"))
mail=json.load(open(args.mail,encoding="utf-8")); doc=Document(args.template)
doc.add_heading("Today's KPI exceptions",level=1)
for month,metric,cite in exceptions[:8]: doc.add_paragraph(f"{month} · {metric} — review required ({cite})",style="List Bullet")
doc.add_heading("Planning decisions and actions",level=1)
for item in mail: doc.add_paragraph(f"{item['summary']} ({item['citation']})",style="List Bullet")
doc.add_heading("Reviewer checklist",level=1)
for item in ["Workbook cells checked","Message citations opened","Recipients and attachments verified","Human approval recorded"]: doc.add_paragraph(item,style="List Bullet")
doc.save(args.output); print(args.output)
'''
    with open(os.path.join(auto, "update_people_workbook.py"), "w", encoding="utf-8") as fh: fh.write(updater)
    with open(os.path.join(auto, "generate_weekly_update.py"), "w", encoding="utf-8") as fh: fh.write(brief)
    with open(os.path.join(inputs, "weekly-input.csv"), "w", newline="", encoding="utf-8") as fh:
        writer = csv.writer(fh); writer.writerow(["Month", "Measure", "Actual"])
        for metric, value in zip(a["case"]["metrics"][:4], [93500, 2, 4, 1]): writer.writerow(["Aug", metric, value])
    findings = [{"summary": "Executive Committee requested the reconciled Q1 milestone list by 4 pm.", "citation": "Outlook message LL-MSG-009"}, {"summary": "Finance owner confirmed the revised contribution assumption needs CFO approval.", "citation": "Outlook message LL-MSG-012"}]
    with open(os.path.join(inputs, "staff-questions.json"), "w", encoding="utf-8") as fh: json.dump(findings, fh, indent=2)
    # The command in the lab expects this convenience filename and template path.
    workbook = os.path.join(folder, f"{stem(a)}-Working-Workbook.xlsx")
    target = os.path.join(folder, "Lumina-Living-People-Tracker.xlsx")
    import shutil; shutil.copy2(workbook, target)
    template = os.path.join(folder, "templates", "Weekly-Update-Template.docx")
    doc = Document(); set_doc_defaults(doc); add_doc_brand(doc, "Daily management brief template")
    doc.add_paragraph("Lumina Living Daily Management Brief", style="Title"); callout(doc, "Control", "Fictional training material. Review every cited cell and message before use.", MINT)
    doc.save(template)
    print("Saved Lab 11 automation starters")


# Each lab ships only the files its own steps name.  A folder full of files the
# lab never opens is the single biggest source of learner confusion.
NEEDS = {
    0:  {"setup"},
    1:  {"workbook"},
    2:  {"workbook"},
    3:  {"brief", "prompt_template"},
    4:  {"brief"},
    5:  {"policyskill"},
    6:  {"workbook"},
    7:  {"brief", "blankdeck"},
    8:  {"workbook"},
    9:  {"hrfolder"},
    10: set(),
    11: {"weekfiles"},
    12: set(),
    13: set(),
    14: {"draftdeck"},
    15: {"workbook", "q1update"},
}


def prune(folder, keep):
    """Delete generated files this lab no longer uses."""
    for root, _dirs, files in os.walk(folder):
        for name in files:
            if name in ("README.md", "TRAINER-GUIDE.md") or name.endswith("-Instructions.pdf"):
                continue
            if name.startswith(".") or name.startswith("~$"):
                os.remove(os.path.join(root, name))
                continue
            if "hr-policy-library" in root or "hr-quarter-files" in root:
                continue
            if os.path.join(root, name) not in keep:
                os.remove(os.path.join(root, name))
                print("Removed unused", os.path.relpath(os.path.join(root, name), folder))


def main():
    for a in ACTS:
        folder = lab_folder(a)
        need = NEEDS.get(a["num"], set())
        before = set()
        for root, _d, files in os.walk(folder):
            before |= {os.path.join(root, f) for f in files}
        if "setup" in need: save_setup_checklist(a, folder)
        if "hrfolder" in need: save_lab09_folder(a, folder)
        if "weekfiles" in need: save_lab10_week(a, folder)
        if "q1update" in need: save_q1_update(a, folder)
        if "q2update" in need: save_q2_update(a, folder)
        if "draftdeck" in need: save_draft_deck(a, folder)
        if "handbook1" in need: save_handbook(a, folder, HANDBOOK_P1, "Staff-Handbook")
        if "handbook2" in need: save_handbook(a, folder, HANDBOOK_P2, "Handbook-Part-2")
        if "policyskill" in need: save_policy_library(a, folder); save_policy_howto(a, folder); save_policy_template(a, folder); save_lab05_source(a, folder)
        if "brief" in need: save_company_brief(a, folder)
        if "sample" in need: save_work_sample(a, folder)
        if "prompt_template" in need: save_prompt_template(a, folder)
        if "workbook" in need: save_workbook(a, folder)
        if "decision_log" in need: save_decision_log(a, folder)
        if "deck" in need: save_presentation(a, folder)
        if "blankdeck" in need: save_blank_deck(a, folder); save_lab07_numbers(a, folder)
        if "automation" in need: save_lab11_automation(a, folder)
        after = set()
        for root, _d, files in os.walk(folder):
            after |= {os.path.join(root, f) for f in files}
        # Anything written or rewritten this run is wanted; the rest is stale.
        made = {p for p in after if p not in before or os.path.getmtime(p) > _START}
        prune(folder, made)


if __name__ == "__main__":
    main()
